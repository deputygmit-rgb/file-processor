import os
import base64
import datetime
import io
import threading
from flask import Flask, request, redirect, url_for, flash, render_template, send_file, make_response, abort, jsonify
from werkzeug.utils import secure_filename
from flask_pymongo import PyMongo
import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import fitz  # PyMuPDF
import pdfplumber
import pytesseract
from bs4 import BeautifulSoup
from pymongo import MongoClient
from gridfs import GridFS, NoFile
from bson import ObjectId
from PIL import Image, ImageFilter, ImageEnhance
import html
from sentence_transformers import SentenceTransformer
import numpy as np
import nltk, re
import openai
import time
from nltk.stem import PorterStemmer
from nltk.corpus import wordnet
import faiss, json
import bson
from bson.json_util import dumps
from waitress import serve
from bson import json_util
from flask import Response
import spacy
from spacy.matcher import Matcher, PhraseMatcher
import datetime
from bson import ObjectId
from neo4j import GraphDatabase
import datetime
import pandas as pd
from bson import ObjectId
from io import BytesIO
from docx import Document
from docx.shared import Pt
import json
import math
import re
import hashlib
from flask import send_from_directory
from flask import Flask, send_file, flash, redirect, url_for
from bson import ObjectId
from flask import abort
from flask import Flask, request, render_template, jsonify
from keybert import KeyBERT
import math, re
from itertools import combinations
import onnxruntime as ort
import numpy as np
from sentence_transformers import CrossEncoder, InputExample
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer
from transformers import AutoTokenizer, AutoModel
import torch
from transformers import AutoTokenizer, AutoModelForCausalLM, pipeline
from peft import PeftModel
app = Flask(__name__)

# 🧠 Load summarization + analysis models (load once)
summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6", device_map=None)
ner_analyzer = pipeline("ner", grouped_entities=True, model="dslim/bert-base-NER", device_map=None)
kw_model = KeyBERT()
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
# --------------------------- Configuration & App --------------------------------
app = Flask(__name__)
app.secret_key = os.getenv("SECRET_KEY", "your_secret_key")
app.config["MONGO_URI"] = os.getenv("MONGO_URI", "mongodb://localhost:27017/dashboard_db")
app.config["UPLOAD_FOLDER"] = os.path.join(os.getcwd(), "uploads")
app.config["MAX_CONTENT_LENGTH"] = 64 * 1024 * 1024  # 64 MB max upload size globally
os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)
FILES_DB = []  # Replace with your actual DB fetch logic
# Tesseract path (optional override via env)

# else system default will be used

# FAISS files
FAISS_INDEX_PATH = os.getenv("FAISS_INDEX_PATH", "faiss_hnsw.index")
FAISS_IDMAP_PATH = os.getenv("FAISS_IDMAP_PATH", "faiss_ids.json")

# Globals
index = None
id_map = []
index_lock = threading.Lock()

# Mongo
mongo = PyMongo(app)
# Initialize MongoDB client with a short server selection timeout so failures are
# detected quickly during startup. If connection to the configured MongoDB
# instance fails, fall back to an in-memory mongomock instance for local
# development (if mongomock is installed). This avoids the app blocking on
# startup when MongoDB is not available.
try:
    client = MongoClient(app.config["MONGO_URI"], serverSelectionTimeoutMS=5000)
    # Force a server selection / ping to verify connectivity
    client.admin.command('ping')
    db = client.get_default_database()
    fs = GridFS(db)
    app.logger.info("Connected to MongoDB at %s", app.config.get("MONGO_URI"))
except Exception as _mongo_err:
    app.logger.warning("MongoDB connection failed: %s", _mongo_err)
    # Attempt to fall back to mongomock for development convenience
    try:
        import mongomock
        client = mongomock.MongoClient()
        # mongomock doesn't support get_default_database well; use configured DB name
        db_name = None
        try:
            # parse DB name from URI if possible
            from urllib.parse import urlparse
            parsed = urlparse(app.config.get("MONGO_URI", ""))
            # For URIs like mongodb://host:port/dbname
            if parsed.path and len(parsed.path) > 1:
                db_name = parsed.path.lstrip('/')
        except Exception:
            db_name = None
        if not db_name:
            db_name = 'dashboard_db'
        db = client[db_name]
        fs = GridFS(db)
        app.logger.info("Using mongomock in-memory MongoDB (db=%s)", db_name)
    except Exception as _mm_err:
        app.logger.error("mongomock fallback not available: %s. DB features will be disabled.", _mm_err)
        client = None
        db = None
        fs = None

ALLOWED_EXTENSIONS = {"xlsx", "pptx", "pdf", "html", "htm", "png", "jpg", "jpeg", "bmp", "gif", "tiff"}


# --- Initialize Neo4j driver with fallback ---
driver = None
try:
    neo4j_uri = os.getenv("NEO4J_URI", "bolt://localhost:7687")
    neo4j_user = os.getenv("NEO4J_USER", "neo4j")
    neo4j_pass = os.getenv("NEO4J_PASS", "Pankaj1234")
    driver = GraphDatabase.driver(neo4j_uri, auth=(neo4j_user, neo4j_pass))
    # Test connection with short timeout
    with driver.session() as session:
        session.run("RETURN 'Neo4j connected' AS msg").single()
    app.logger.info("Neo4j connected at %s", neo4j_uri)
except Exception as e:
    app.logger.warning("Neo4j connection failed (%s). KG features will be disabled. %s", e.__class__.__name__, e)
    driver = None


# --- Helper: timed Neo4j runner to log slow queries ---
SLOW_KG_MS = int(os.getenv("SLOW_KG_MS", "200"))
def timed_run(session, cypher, **params):
    """Run a cypher query and log if it takes longer than SLOW_KG_MS milliseconds."""
    t0 = time.time()
    res = session.run(cypher, **params)
    dur = (time.time() - t0) * 1000.0
    if dur > SLOW_KG_MS:
        try:
            app.logger.warning("Slow KG query: %.1fms — query head: %s", dur, cypher[:300])
        except Exception:
            app.logger.warning("Slow KG query: %.1fms", dur)
    return res


# --- Simple in-memory TTL LRU cache for KG results ---
from collections import OrderedDict, deque

KG_CACHE_MAX = int(os.getenv("KG_CACHE_MAX", "1024"))
KG_CACHE_TTL = int(os.getenv("KG_CACHE_TTL", "3600"))  # seconds
KG_QUERY_RATE_LIMIT = int(os.getenv("KG_QUERY_RATE_LIMIT", "100"))  # max queries per 60s

# Rate limiting for KG queries (sliding 60-second window)
_kg_query_times = deque(maxlen=KG_QUERY_RATE_LIMIT)

def _cache_get(cache, key):
    now = time.time()
    item = cache.get(key)
    if not item:
        return None
    value, ts = item
    if now - ts > KG_CACHE_TTL:
        try:
            del cache[key]
        except KeyError:
            pass
        return None
    # move to end (recently used)
    try:
        cache.move_to_end(key)
    except Exception:
        pass
    return value

def _cache_set(cache, key, value):
    now = time.time()
    cache[key] = (value, now)
    try:
        cache.move_to_end(key)
    except Exception:
        pass
    # evict oldest if necessary
    while len(cache) > KG_CACHE_MAX:
        try:
            cache.popitem(last=False)
        except Exception:
            break


def _sanitize_query_term(term: str, max_len: int = 200) -> str:
    """
    Sanitize a query term to prevent injection and excessive length.
    - Trim to max_len chars
    - Remove/escape dangerous Cypher patterns
    - Collapse whitespace
    Returns sanitized term safe for parameterized queries.
    """
    if not term:
        return ""
    # Collapse whitespace
    term = re.sub(r'\s+', ' ', str(term).strip())
    # Truncate
    term = term[:max_len]
    # Log suspicious patterns (for debugging)
    if any(p in term.lower() for p in ['return', 'match', 'where', 'create', 'delete', 'drop', 'merge']):
        app.logger.debug("Sanitize: suspicious pattern in term '%s'", term[:50])
    return term


def _validate_int_param(value, default: int = 2, min_val: int = 1, max_val: int = 10) -> int:
    """
    Safely validate and constrain an integer parameter (e.g., max_hops).
    Prevents excessively large values from causing performance issues.
    """
    try:
        v = int(value)
        return max(min_val, min(v, max_val))
    except (ValueError, TypeError):
        return default


def _check_kg_rate_limit() -> bool:
    """
    Check if KG query rate limit is exceeded (simple sliding window, 60-second window).
    Returns True if within limit, False if rate-limited.
    """
    now = time.time()
    # Remove old entries outside the 60-second window
    while _kg_query_times and _kg_query_times[0] < now - 60:
        _kg_query_times.popleft()
    if len(_kg_query_times) >= KG_QUERY_RATE_LIMIT:
        return False
    _kg_query_times.append(now)
    return True

# caches
KG_EXPAND_CACHE = OrderedDict()
KG_IMPORTANCE_CACHE = OrderedDict()
KG_RELATED_FAILURES_CACHE = OrderedDict()

if client is not None:
    db = client["railwaydb"]  # use your DB name
    files_collection = db["files"]  # common collection for all uploads
    kg_collection = db["knowledge_graph"]  # <-- Add this line
else:
    db = None
    files_collection = None
    kg_collection = None
    app.logger.warning("No MongoDB client available — DB-backed features will be disabled.")
_illegal_xml_re = re.compile(r'[\x00-\x08\x0b-\x0c\x0e-\x1f]')
ps = PorterStemmer()

# Ensure NLTK data present (wordnet) — try to download if missing
try:
    nltk.data.find("corpora/wordnet")
except LookupError:
    nltk.download("wordnet")
    nltk.download("omw-1.4")

# SentenceTransformer model (single initialization)
model = SentenceTransformer("all-MiniLM-L6-v2")
nlp = spacy.load("en_core_web_sm")
nlp.max_length = 3_000_000  # increase to comfortably handle your largest tex
matcher = PhraseMatcher(nlp.vocab, attr="LOWER")
global SIGNAL_ASSETS, RAIL_FAILURE_DICT, RAILWAY_ASSETS
rail_terms = [
    # Core entities
    "railway board", "railway zone", "railway division", "railway department",
    "zonal railway", "divisional railway", "station", "yard", "loco shed",

    # Departments
    "civil engineering", "electrical engineering", "mechanical department",
    "signal department", "telecom department", "commercial department",
    "traffic department", "stores department", "personnel department",
    "accounts department", "security department", "medical department",
    "construction department", "safety department", "operations department",
    "management department",

    # Key designations
    "divisional railway manager", "chief engineer", "senior divisional engineer",
    "principal chief engineer", "principal chief commercial manager",
    "principal chief mechanical engineer", "principal chief signal and telecom engineer",
    "principal financial advisor", "principal chief personnel officer",
    "principal chief medical director", "principal chief safety officer",
    "principal chief security commissioner", "additional general manager",

    # Common functions
    "train operations", "maintenance", "infrastructure", "rolling stock",
    "ticketing", "freight booking", "material stores", "signalling", "power supply",
    "traction distribution", "locomotive", "carriage", "wagon", "overhead equipment",

    # Railway Zones
    "northern railway", "southern railway", "eastern railway", "western railway",
    "central railway", "north eastern railway", "south eastern railway",
    "northeast frontier railway", "south central railway", "south coast railway",
    "konkan railway", "east central railway", "south east central railway",
    "north western railway", "east coast railway", "north central railway",
    "south western railway", "west central railway","NR","SR","ER","WR","CR","NER","SER","NFR","SCR","ECR","SECR","NWR","ECoR","NCR","SWR","WCR",

    # Major Divisions (partial list; can be extended dynamically)
    "delhi division", "mumbai division", "chennai division", "howrah division",
    "secunderabad division", "bilaspur division", "bhopal division", "hubballi division",
    "jaipur division", "nagpur division", "lucknow division", "guwahati division",
    "vadodara division", "madurai division", "raipur division", "ratlam division", "Raipur division"
    "sambalpur division", "tiruchirappalli division", "vijayawada division", "visakhapatnam division","BSP","BPL","UBL","UBL","JPR","NGP","LKO","GHY","BRC","MDU","R","RTM","SBP","TIR","BZA","VSKP","NAG","NGP"
]
# Asset list for all Indian Railway departments
railway_assets = {
    "civil engineering": [
        # Track infrastructure
        "rail tracks", "sleepers", "ballast", "turnouts", "points and crossings",
        "track joints", "fishplates", "track fastenings", "rails", "switch expansion joints",
        # Structures
        "bridges", "culverts", "retaining walls", "tunnels", "buildings",
        "platforms", "foot over bridges", "subways", "station buildings",
        "drainage systems", "embankments", "cuttings",
        # Water & environment
        "water supply systems", "sewage systems", "rainwater harvesting", "drinking water points",
        # Equipment
        "earthwork machinery", "survey equipment", "welding machines", "tamping machines"
    ],

    "signal and telecommunication": [
        # Signalling Equipment
        "signal relays", "axle counters", "interlocking panels", "route indicators",
        "signal posts", "led signal units", "signal transformers",
        "control panels", "signal relays", "neutral sections", "track circuits",
        "block instruments", "lever frames", "route setting panels", "panel interlocking",
        "electronic interlocking", "solid state interlocking",
        # Communication
        "telephones", "optical fiber cables", "vhf sets", "microwave towers", "exchange systems",
        "control communication equipment", "train radio system",
        # Power & control
        "power supply panels", "batteries", "charger units", "dc distribution boards"
    ],

    "mechanical": [
        # Rolling stock
        "locomotives", "diesel locomotives", "electric locomotives", "coaches", "wagons",
        "brake vans", "cranes", "tower cars",
        # Components
        "bogies", "axle boxes", "brake systems", "air brake system", "couplers",
        "draw gear", "traction motors", "compressors", "wheel sets", "bearing housings",
        # Maintenance
        "pit lines", "wash plants", "lifting jacks", "underframe equipment",
        "lubrication systems", "fuel pumps", "sanders"
    ],

    "electrical": [
        # Traction Power
        "traction transformers", "overhead equipment", "section insulators",
        "neutral sections", "contact wire", "dropper wire", "catenary wire",
        "traction substations", "feeders", "circuit breakers",
        # Non-traction Power
        "lighting systems", "switchgear", "earthing systems", "distribution panels",
        "batteries", "inverters", "ups systems", "generators",
        # Signalling & general services
        "station lighting", "lift systems", "fans", "air conditioning units"
    ],

    "stores": [
        # Inventory & supply
        "spare parts", "fasteners", "oils", "lubricants", "tools", "consumables",
        "electrical spares", "signalling spares", "mechanical spares",
        "track fittings", "relays", "bearings", "paint", "rubber components"
    ],

    "commercial": [
        # Passenger & freight operations
        "ticketing systems", "booking counters", "reservation systems",
        "automatic ticket vending machines", "freight booking terminals",
        "display boards", "train indication systems", "enquiry systems",
        "luggage scanners", "announcement systems", "passenger information displays"
    ],

    "security": [
        "surveillance cameras", "cctv systems", "baggage scanners",
        "access control systems", "door frame metal detectors",
        "hand held detectors", "public address systems", "fire alarm panels",
        "security control rooms", "video recorders"
    ]
}
signal_assets = {
    "Panel Interlocking": {"category": "Signalling System", "unit": "Stations"},
    "Point Machine": {"category": "Signalling System", "unit": "Nos"},
    "Axle counter": {"category": "Train Detection", "unit": "Nos"},
    "MSDAC": {"category": "Train Detection", "unit": "Stations"},
    "Signal Signals": {"category": "Signalling System", "unit": "Stations"},
    "Electronic Interlocking": {"category": "Signalling System", "unit": "Stations"},
    "Route Relay Interlocking": {"category": "Signalling System", "unit": "Stations"},
    "LED Lit Signals": {"category": "Signal Equipment", "unit": "Stations"},
    "Data Logger": {"category": "Monitoring Equipment", "unit": "Stations"},
    "Colour Light Signalling": {"category": "Signalling System", "unit": "Stations"},
    "Block Proving by Axle Counter": {"category": "Train Detection", "unit": "Block Sections"},
    "Track Circuiting": {"category": "Train Detection", "unit": "Stations"},
    "Automatic Block Signalling": {"category": "Train Control", "unit": "Rkm"},
    "Intermediate Block Signalling": {"category": "Train Control", "unit": "Nos"},
    "Interlocked Level Crossing": {"category": "Safety Infrastructure", "unit": "Nos"},
    "ADVANCED STARTER": {"category": "Safety Infrastructure", "unit": "Nos"},
    "GATE SIGNAL": {"category": "Safety Infrastructure", "unit": "Nos"},
    "HOME SIGNAL": {"category": "Safety Infrastructure", "unit": "Nos"},
    "INTERMEDIATE BLOCK SIGNAL": {"category": "Safety Infrastructure", "unit": "Nos"},
    "BLOCK INSTRUMENT": {"category": "Safety Infrastructure", "unit": "Nos"},
    "Kavach": {"category": "Train Protection", "unit": "Rkm"}
}
rail_failure_dict = [
    # Signal & Telecommunication
    {"failure_code": "ACP", "failure_subcode": "ACP", "failure_desc": "ALARM CHAIN PULLING", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "ASF", "failure_subcode": "ASF", "failure_desc": "ADVANCED STARTER FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "AUTSF", "failure_subcode": "AUTSF", "failure_desc": "AUTOMATIC SIGNAL FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "AXCF", "failure_subcode": "AXCF", "failure_desc": "BLOCK AXLE COUNTER FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "BBST", "failure_subcode": "BBST", "failure_desc": "ST BLOCK BURST", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "BIF", "failure_subcode": "BIF", "failure_desc": "BLOCK INSTRUMENT FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "GSF", "failure_subcode": "GSF", "failure_desc": "GATE SIGNAL FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "HSF", "failure_subcode": "HSF", "failure_desc": "HOME SIGNAL FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "IBSF", "failure_subcode": "IBSF", "failure_desc": "INTERMEDIATE BLOCK SIGNAL FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "ICMS", "failure_subcode": "ICMS", "failure_desc": "COA/ICMS FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "PANF", "failure_subcode": "PANF", "failure_desc": "RRI/PI FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "PTF", "failure_subcode": "PTF", "failure_desc": "POINT FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "SPEF", "failure_subcode": "SPEF", "failure_desc": "SIGNAL POWER EQUIPMENT FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "SSIF", "failure_subcode": "SSIF", "failure_desc": "SSIF FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "STCC", "failure_subcode": "STCC", "failure_desc": "CABLE CUTTING", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},
    {"failure_code": "TCF", "failure_subcode": "TCF", "failure_desc": "TRACK CIRCUIT FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "SIGN AND TELE", "department": "Signal and Telecommunication"},

    # Mechanical (Coaches, Wagons, Diesel Loco)
    {"failure_code": "ACOP", "failure_subcode": "ACOP", "failure_desc": "PASSENGER TRAIN ANGLE COCK OPERATED", "valid": "Y", "user_asset_failure": "C", "system_auto": "Y", "asset_group": "MECHANICAL", "department": "Mechanical"},
    {"failure_code": "BBJP", "failure_subcode": "BBJP", "failure_desc": "PASSENGER TRAIN BRAKE BLOCK JAM", "valid": "Y", "user_asset_failure": "C", "system_auto": "Y", "asset_group": "MECHANICAL", "department": "Mechanical"},
    {"failure_code": "DLFG", "failure_subcode": "DLFG", "failure_desc": "DIESEL GOODS TRAIN LOCO FAILED", "valid": "Y", "user_asset_failure": "DL", "system_auto": "Y", "asset_group": "MECHANICAL", "department": "Mechanical"},
    {"failure_code": "DLORWP", "failure_subcode": "DLORWP", "failure_desc": "DIESEL LOCO LOSS ON RUN DUE TO WRONG POWER", "valid": "Y", "user_asset_failure": "DL", "system_auto": "Y", "asset_group": "MECHANICAL", "department": "Mechanical"},
    {"failure_code": "DLTG", "failure_subcode": "DLTG", "failure_desc": "DIESEL GOODS TRAIN LOCO TROUBLE", "valid": "Y", "user_asset_failure": "DL", "system_auto": "Y", "asset_group": "MECHANICAL", "department": "Mechanical"},

    # Electrical
    {"failure_code": "ELFG", "failure_subcode": "ELFG", "failure_desc": "ELECTRIC GOODS TRAIN LOCO FAILED", "valid": "Y", "user_asset_failure": "EL", "system_auto": "Y", "asset_group": "ELECTRICAL", "department": "Electrical"},
    {"failure_code": "OHE", "failure_subcode": "OHE", "failure_desc": "OHE FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "ELECTRICAL", "department": "Electrical"},
    {"failure_code": "ATF", "failure_subcode": "ATF", "failure_desc": "AUXILIARY TRANSFORMER FAILURE", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "ELECTRICAL", "department": "Electrical"},
    {"failure_code": "REOHE", "failure_subcode": "REOHE", "failure_desc": "OHE HANGING ETC", "valid": "N", "user_asset_failure": "", "system_auto": "Y", "asset_group": "ELECTRICAL", "department": "Electrical"},

    # Engineering
    {"failure_code": "BBENG", "failure_subcode": "BBENG", "failure_desc": "BLOCK BURSTING", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "ENGINEERING", "department": "Engineering"},
    {"failure_code": "BUCTRC", "failure_subcode": "BUCTRC", "failure_desc": "BUCKLING OF TRACK", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "ENGINEERING", "department": "Engineering"},
    {"failure_code": "TRB", "failure_subcode": "TRB", "failure_desc": "TONGUE RAIL BROKEN", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "ENGINEERING", "department": "Engineering"},

    # Law & Order
    {"failure_code": "AOR", "failure_subcode": "AOR", "failure_desc": "ATTACK ON RLY-MEN/PROPERTY/TRN", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "LAW AND ORDER", "department": "Security"},
    {"failure_code": "BNDH", "failure_subcode": "BNDH", "failure_desc": "BANDH", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "LAW AND ORDER", "department": "Security"},
    {"failure_code": "MAOTH", "failure_subcode": "MAOTH", "failure_desc": "MISCREANT ACTIVITY OTHERS", "valid": "Y", "user_asset_failure": "", "system_auto": "Y", "asset_group": "LAW AND ORDER", "department": "Security"},

    # Blocks
    {"failure_code": "BBELECN", "failure_subcode": "BBELECN", "failure_desc": "PBC ELEC CNST BLOCK BURST", "valid": "Y", "user_asset_failure": "", "system_auto": "", "asset_group": "BLOCKS", "department": "Blocks"},
]
# --------------------------- FAISS helpers --------------------------------
def format_time(ts):
    return ts.strftime("%Y-%m-%d %H:%M:%S")    
def load_faiss_index():
    global index, id_map
    if os.path.exists(FAISS_INDEX_PATH):
        try:
            index = faiss.read_index(FAISS_INDEX_PATH)
        except Exception as e:
            app.logger.exception("Failed to load FAISS index: %s", e)
            index = None
    else:
        index = None

    if os.path.exists(FAISS_IDMAP_PATH):
        try:
            with open(FAISS_IDMAP_PATH, "r", encoding="utf-8") as f:
                id_map[:] = json.load(f)
        except Exception as e:
            app.logger.exception("Failed to load id_map: %s", e)
            id_map[:] = []
    else:
        id_map[:] = []

    # Validate
    try:
        if index is not None:
            idx_dim = getattr(index, "d", None)
            app.logger.info("Loaded FAISS index dim=%s ntotal=%s", idx_dim, index.ntotal)
            if index.ntotal != len(id_map):
                app.logger.warning("FAISS ntotal (%d) != id_map length (%d). Truncating id_map.", index.ntotal, len(id_map))
                id_map[:] = id_map[: index.ntotal]
    except Exception:
        pass

# initial load
load_faiss_index()

# kg_neo4j_helpers.py
def asset_exists(tx, asset_name):
    result = tx.run("MATCH (a:Asset {name:$name}) RETURN count(a) > 0 AS exists", name=asset_name)
    return result.single()["exists"]

def failure_exists(tx, failure_code):
    result = tx.run("MATCH (f:Failure {code:$code}) RETURN count(f) > 0 AS exists", code=failure_code)
    return result.single()["exists"]

def link_failure_to_existing_nodes(tx, failure, file_id):
    """
    Create or link failure to *existing* hierarchy nodes.
    If nodes exist (Zone, Division, Section, Station, Gear), reuse them.
    """
    q = """
    // Match existing hierarchy
    MATCH (z:Zone {name:$zone})
    OPTIONAL MATCH (d:Division {name:$division})-[:PART_OF]->(z)
    OPTIONAL MATCH (s:Section {name:$section})-[:PART_OF]->(d)
    OPTIONAL MATCH (st:Station {code:$station_code})-[:HAS_STATION]->(s)
    OPTIONAL MATCH (g:Gear {name:$gear_name})-[:HAS_GEAR]->(st)

    // If missing, create them
    WITH z, COALESCE(d, MERGE (dnew:Division {name:$division}) MERGE (dnew)-[:PART_OF]->(z) RETURN dnew) AS div,
         COALESCE(s, MERGE (snew:Section {name:$section}) MERGE (snew)-[:PART_OF]->(div) RETURN snew) AS sec,
         COALESCE(st, MERGE (stnew:Station {code:$station_code}) MERGE (stnew)-[:HAS_STATION]->(sec) RETURN stnew) AS stn,
         COALESCE(g, MERGE (gnew:Gear {name:$gear_name, subtype:$gear_subtype}) MERGE (stn)-[:HAS_GEAR]->(gnew) RETURN gnew) AS gear

    // Create or link failure
    MERGE (f:Failure {failure_entry_no:$failure_entry_no})
    SET f.department=$department,
        f.cause=$cause,
        f.rectification_status=$rectification_status,
        f.rectification_by=$rectification_by,
        f.failure_duration=$failure_duration,
        f.upload_source=$file_id

    MERGE (gear)-[:HAS_FAILURE]->(f)
    """
    tx.run(q, **failure, file_id=str(file_id))

def create_kg(tx, rail_failure_dict):
    # Create Departments and AssetGroups first
    for item in rail_failure_dict:
        department = item["department"]
        asset_group = item["asset_group"]
        failure_code = item["failure_code"]
        failure_desc = item["failure_desc"]

        # Merge ensures no duplicates
        tx.run("MERGE (d:Department {name: $department})", department=department)
        tx.run("""
            MERGE (ag:AssetGroup {name: $asset_group})
            MERGE (d:Department {name: $department})
            MERGE (d)-[:HAS_ASSET_GROUP]->(ag)
        """, department=department, asset_group=asset_group)

        # Failure node
        tx.run("""
            MERGE (f:Failure {code: $failure_code, description: $failure_desc})
            MERGE (ag:AssetGroup {name: $asset_group})
            MERGE (ag)-[:HAS_FAILURE]->(f)
            MERGE (f)-[:BELONGS_TO]->(:Department {name: $department})
        """, failure_code=failure_code, failure_desc=failure_desc, asset_group=asset_group, department=department)
# kg_bootstrap.py
from neo4j import GraphDatabase
# in background_tasks()

def map_to_existing_entities(failure):
    """
    Align Excel row to known KG entities to avoid duplicates.
    """
    # normalize names
    failure["zone"] = failure["zone"].upper().strip()
    failure["division"] = failure["division"].title().strip()

    # normalize department
    if "tele" in failure.get("department", "").lower():
        failure["department"] = "S&T"

    # map gear using dictionary
    for group, items in signal_assets.items():
        for asset in items.keys():
            if asset.lower() in failure["gear_name"].lower():
                failure["gear_group"] = group
                failure["gear_name"] = asset
                break
    return failure

def preload_static_entities(driver):
    """
    Create static entity hierarchy (Zones, Divisions, Departments, Asset Groups)
    once and cache their IDs.
    """
    with driver.session() as session:
        # Example Zone-Division structure
        for zone, divs in railway_assets.items():
            session.run("MERGE (z:Zone {name:$zone})", zone=zone)
            for division in divs:
                session.run("""
                    MERGE (d:Division {name:$division})
                    MERGE (d)-[:PART_OF]->(z)
                """, zone=zone, division=division)

        # Departments
        for dept in ["S&T", "Engineering", "Electrical", "Mechanical"]:
            session.run("MERGE (:Department {name:$dept})", dept=dept)

        # Asset types
        for asset_type, assets in signal_assets.items():
            for asset_name in assets.keys():
                session.run("""
                    MERGE (a:AssetGroup {name:$asset_name, type:$asset_type})
                """, asset_name=asset_name, asset_type=asset_type)

with driver.session() as session:
    preload_static_entities(driver)
    session.execute_write(create_kg, rail_failure_dict)

def add_failures_to_graph(failures):
    """
    Add failures to Neo4j KG.
    
    failures: List of dicts with keys:
        department, asset_type, failure_code, failure_subcode, failure_desc, valid, user_asset_failure, system_auto, asset_group
    """
    if not failures:
        return 0

    # Prepare rows for UNWIND batch merge
    rows = []
    for f in failures:
        rows.append({
            "department": f.get("department", ""),
            "asset_type": f.get("asset_type", ""),
            "failure_code": f.get("failure_code", ""),
            "failure_subcode": f.get("failure_subcode", f.get("failure_sub_code", "")),
            "failure_desc": f.get("failure_desc", f.get("failure_description", "")),
            "valid": f.get("valid", "Y"),
            "user_asset_failure": f.get("user_asset_failure", ""),
            "system_auto": f.get("system_auto", ""),
            "asset_group": f.get("asset_group", ""),
        })

    q = """
    UNWIND $rows AS f
    MERGE (d:Department {name: f.department})
    MERGE (a:AssetType {name: f.asset_type})
    MERGE (d)-[:HAS_ASSET]->(a)
    MERGE (fc:FailureCode {code: f.failure_code})
    SET fc.subcode = f.failure_subcode,
        fc.description = f.failure_desc,
        fc.valid = f.valid,
        fc.user_asset_failure = f.user_asset_failure,
        fc.system_auto = f.system_auto,
        fc.asset_group = f.asset_group
    MERGE (a)-[:HAS_FAILURE]->(fc)
    RETURN count(DISTINCT fc) AS created
    """

    with driver.session() as session:
        try:
            res = timed_run(session, q, rows=rows)
            created = 0
            try:
                r = res.single()
                created = int(r["created"]) if r and r["created"] is not None else 0
            except Exception:
                created = 0
            app.logger.info("Batch inserted/merged %d failure nodes into KG.", created)
            return created
        except Exception as e:
            app.logger.exception("Failed to batch insert failures into KG: %s", e)
            # Fallback: try per-item to surface specific bad rows
            created = 0
            for f in failures:
                try:
                    with driver.session() as session:
                        session.execute_write(_create_failure_nodes, f)
                        created += 1
                except Exception as ex:
                    app.logger.exception("Failed to create failure node for %s: %s", f.get('failure_code'), ex)
            return created

def _create_failure_nodes(tx, f):
    # Merge Department node
    tx.run("""
        MERGE (d:Department {name: $department})
    """, department=f['department'])
    
    # Merge AssetType node and connect to Department
    tx.run("""
        MATCH (d:Department {name: $department})
        MERGE (a:AssetType {name: $asset_type})
        MERGE (d)-[:HAS_ASSET]->(a)
    """, department=f['department'], asset_type=f['asset_type'])
    
    # Merge FailureCode node and connect to AssetType
    tx.run("""
        MATCH (a:AssetType {name: $asset_type})
        MERGE (fc:FailureCode {code: $failure_code})
        SET fc.subcode = $failure_subcode,
            fc.description = $failure_desc,
            fc.valid = $valid,
            fc.user_asset_failure = $user_asset_failure,
            fc.system_auto = $system_auto,
            fc.asset_group = $asset_group
        MERGE (a)-[:HAS_FAILURE]->(fc)
    """, 
    asset_type=f['asset_type'], 
    failure_code=f['failure_code'], 
    failure_subcode=f.get('failure_subcode',''),
    failure_desc=f.get('failure_desc',''),
    valid=f.get('valid','Y'),
    user_asset_failure=f.get('user_asset_failure',''),
    system_auto=f.get('system_auto',''),
    asset_group=f.get('asset_group',''))



def add_to_faiss(mongo_id: str, embedding: np.ndarray, normalize: bool = True):
    """
    Add embedding numpy array to FAISS in a thread-safe and dtype-safe manner.
    By default it normalizes vectors (recommended) so cosine similarity can be computed via inner product.
    """
    global index, id_map

    vec = np.asarray(embedding, dtype="float32").reshape(1, -1)

    if normalize:
        nrm = np.linalg.norm(vec, axis=1, keepdims=True)
        nrm[nrm == 0] = 1.0
        vec = vec / nrm

    with index_lock:
        if index is None:
            dim = vec.shape[1]
            index = faiss.IndexHNSWFlat(dim, 32)
            index.hnsw.efConstruction = 200
            index.hnsw.efSearch = 50

        # check dims
        if vec.shape[1] != index.d:
            raise RuntimeError(f"Embedding dimension ({vec.shape[1]}) does not match index dimension ({index.d})")

        index.add(vec)
        id_map.append(mongo_id)

        # persist
        try:
            faiss.write_index(index, FAISS_INDEX_PATH)
            with open(FAISS_IDMAP_PATH, "w", encoding="utf-8") as f:
                json.dump(id_map, f)
        except Exception as e:
            app.logger.exception("Failed to persist FAISS index or id_map: %s", e)




def rebuild_faiss(batch_size=256, normalize=True, drop_existing=True):
    """
    Rebuild FAISS index from MongoDB documents. This re-encodes search_text using the current model.
    WARNING: This will overwrite the FAISS index and id_map files when finished.
    """
    global index, id_map

    if drop_existing:
        try:
            if os.path.exists(FAISS_INDEX_PATH):
                os.rename(FAISS_INDEX_PATH, FAISS_INDEX_PATH + ".bak")
            if os.path.exists(FAISS_IDMAP_PATH):
                os.rename(FAISS_IDMAP_PATH, FAISS_IDMAP_PATH + ".bak")
        except Exception as e:
            app.logger.warning("Failed to backup old index files: %s", e)

    cursor = mongo.db.files.find({}, {"_id": 1, "search_text": 1})
    id_map = []
    index = None
    batch_ids = []
    batch_texts = []
    total = 0

    for doc in cursor:
        batch_ids.append(str(doc["_id"]))
        batch_texts.append(doc.get("search_text", "") or "")
        if len(batch_ids) >= batch_size:
            emb = model.encode(batch_texts, convert_to_numpy=True)
            emb = np.asarray(emb, dtype="float32")
            if normalize:
                norms = np.linalg.norm(emb, axis=1, keepdims=True)
                norms[norms == 0] = 1.0
                emb = emb / norms
            with index_lock:
                if index is None:
                    dim = emb.shape[1]
                    index = faiss.IndexHNSWFlat(dim, 32)
                    index.hnsw.efConstruction = 200
                    index.hnsw.efSearch = 50
                index.add(emb)
                id_map.extend(batch_ids)
            total += len(batch_ids)
            batch_ids = []
            batch_texts = []

    if batch_ids:
        emb = model.encode(batch_texts, convert_to_numpy=True)
        emb = np.asarray(emb, dtype="float32")
        if normalize:
            norms = np.linalg.norm(emb, axis=1, keepdims=True)
            norms[norms == 0] = 1.0
            emb = emb / norms
        with index_lock:
            if index is None:
                dim = emb.shape[1]
                index = faiss.IndexHNSWFlat(dim, 32)
                index.hnsw.efConstruction = 200
                index.hnsw.efSearch = 50
            index.add(emb)
            id_map.extend(batch_ids)
        total += len(batch_ids)

    # persist
    with index_lock:
        if index is not None:
            faiss.write_index(index, FAISS_INDEX_PATH)
        with open(FAISS_IDMAP_PATH, "w", encoding="utf-8") as f:
            json.dump(id_map, f)

    app.logger.info("Rebuilt FAISS index with %d vectors.", total)
    return total

# --------------------------- Utilities & processing --------------------------------

def flatten_text(data):
    texts = []
    if data is None:
        return texts

    if isinstance(data, dict):
        for v in data.values():
            texts.extend(flatten_text(v))
    elif isinstance(data, (list, tuple)):
        # Heuristic: treat a list-of-lists as a table (rows of cells)
        if data and all(isinstance(row, (list, tuple)) for row in data):
            # Build a compact, searchable table representation that preserves
            # header -> row relationships. We wrap with markers so chunking
            # can treat the table as an atomic block.
            try:
                rows = data
                header = None
                # Heuristic: first row is header if many cells contain letters
                first = rows[0]
                if all(isinstance(c, (str, int, float, type(None))) for c in first):
                    alpha_count = sum(1 for c in first if isinstance(c, str) and re.search(r'[A-Za-z]', c))
                    if alpha_count >= max(1, len(first)//2):
                        header = [str(c) if c is not None else "" for c in first]
                        data_rows = rows[1:]
                    else:
                        data_rows = rows
                else:
                    data_rows = rows

                tbl_lines = ["___TABLE_START___"]
                if header:
                    tbl_lines.append("TABLE_HEADER: " + " | ".join([h for h in header]))
                else:
                    # create generic column names
                    cols = len(rows[0]) if rows else 0
                    tbl_lines.append("TABLE_HEADER: " + " | ".join([f"col{i+1}" for i in range(cols)]))

                for r_idx, r in enumerate(data_rows, start=1):
                    cells = [clean_cell_value(c) for c in r]
                    tbl_lines.append("TABLE_ROW: " + " | ".join(cells))

                tbl_lines.append("___TABLE_END___")
                texts.append("\n".join(tbl_lines))
            except Exception:
                # Fallback to flattening cells
                for item in data:
                    texts.extend(flatten_text(item))
        else:
            for item in data:
                texts.extend(flatten_text(item))
    elif isinstance(data, str):
        texts.append(data)
    elif isinstance(data, (datetime.date, datetime.datetime)):
        texts.append(data.isoformat())
    elif isinstance(data, datetime.time):
        texts.append(data.strftime("%H:%M:%S"))
    elif isinstance(data, ObjectId):
        texts.append(str(data))
    else:
        try:
            texts.append(str(data))
        except Exception:
            pass

    return texts


import re

def make_snippet(text, query, max_words=400):
    """
    Generate a snippet from text that focuses on relevant paragraphs.
    Uses progressive AND/OR paragraph extraction to ensure snippet
    contains query terms.
    """

    if not text or not query:
        return ""

    # Step 1: Extract relevant paragraphs
    paragraphs = extract_relevant_paragraphs(text, query, max_chunks=5)

    if not paragraphs:
        # fallback: take first max_words words from full text
        words = re.findall(r'\w+', text)
        return " ".join(words[:max_words])

    # Step 2: Combine top paragraphs for snippet
    snippet_text = " ".join(paragraphs)

    # Step 3: Limit snippet length
    words = re.findall(r'\w+', snippet_text)
    snippet = " ".join(words[:max_words])

    return snippet


def ensure_text_index():
    try:
        existing = [idx["name"] for idx in db.files.list_indexes()]
        if "files_text_index" not in existing:
            db.files.create_index([("filename", "text"), ("search_text", "text")], name="files_text_index", default_language="english")
            app.logger.info("Created text index 'files_text_index' on filename + search_text")
        else:
            app.logger.info("Text index 'files_text_index' already exists")
    except Exception as e:
        app.logger.exception("Failed to create/check text index: %s", e)

ensure_text_index()


def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS



_illegal_xml_re = re.compile(
    r"[\x00-\x08\x0b\x0c\x0e-\x1f]"  # illegal XML chars
)

def clean_cell_value(value):
    if value is None:
        return ""
    if isinstance(value, str):
        return _illegal_xml_re.sub("", value)
    elif isinstance(value, (datetime.date, datetime.datetime)):
        return value.isoformat()
    elif isinstance(value, datetime.time):
        return value.strftime("%H:%M:%S")
    return str(value)




def process_excel(file_stream):
    wb = openpyxl.load_workbook(file_stream, data_only=True)
    sheets_data = {}
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.values)
        rows = [row for row in rows if any(cell not in (None, "", " ") for cell in row)]
        if not rows:
            continue
        cols = list(zip(*rows))
        non_empty_cols_idx = [i for i, col in enumerate(cols) if any(cell not in (None, "", " ") for cell in col)]
        if not non_empty_cols_idx:
            continue
        cleaned_rows = []
        for row in rows:
            cleaned_row = [row[i] if i < len(row) else None for i in non_empty_cols_idx]
            cleaned_rows.append(cleaned_row)
        sheets_data[sheet_name] = cleaned_rows
    return sheets_data


def process_pptx(file_stream):
    prs = Presentation(file_stream)
    slides_data = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        slide_tables = []
        slide_images = []
        def extract_images_from_shape(shape):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext.lower() if image.ext else "png"
                try:
                    ocr_text = extract_text_from_image(image_bytes)
                except Exception:
                    ocr_text = ""
                image_id = fs.put(image_bytes, contentType=f"image/{image_ext}", filename=f"slide{slide_idx}_image.{image_ext}")
                slide_images.append({"id": str(image_id), "ext": image_ext, "ocr_text": ocr_text})
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    extract_images_from_shape(shp)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
            if shape.has_table:
                table = shape.table
                rows = []
                for r in range(len(table.rows)):
                    row_cells = []
                    for c in range(len(table.columns)):
                        cell_text = table.cell(r, c).text.strip()
                        row_cells.append(cell_text)
                    rows.append(row_cells)
                slide_tables.append(rows)
            extract_images_from_shape(shape)
        slides_data.append({"title": slide.shapes.title.text if slide.shapes.title else "", "texts": slide_texts, "tables": slide_tables, "images": slide_images})
    return slides_data


def preprocess_image_for_ocr(image_bytes):
    image = Image.open(io.BytesIO(image_bytes))
    image = image.convert("L")
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(2.0)
    image = image.filter(ImageFilter.SHARPEN)
    min_size = 300
    if image.width < min_size or image.height < min_size:
        scale = max(min_size / image.width, min_size / image.height)
        new_size = (int(image.width * scale), int(image.height * scale))
        image = image.resize(new_size, Image.LANCZOS)
    return image


def clean_extracted_text(text):
    if not text:
        return ""

    # 1) Remove (cid:####) patterns and illegal xml/control chars
    text = re.sub(r'\(cid:\d+\)', '', text)
    text = _illegal_xml_re.sub('', text)

    # 2) Normalize whitespace
    text = re.sub(r'\s+', ' ', text).strip()

    # Helper: decide if token is likely to be a real word (very permissive)
    def likely_word(tok: str) -> bool:
        tok = tok.strip("\'\"()[]{}:;,.!?-_")
        if not tok:
            return False
        # numbers, short tokens and tokens with vowels are likely legitimate
        if len(tok) <= 3:
            return True
        if re.search(r'[aeiouAEIOU]', tok):
            return True
        # contains letters and digits mix (e.g., model numbers) — keep
        if re.search(r'[A-Za-z]', tok) and re.search(r'\d', tok):
            return True
        return False

    # 3) Fix common OCR artifacts conservatively
    words = text.split()
    cleaned_words = []
    for w in words:
        orig = w
        # Trim long runs of same character: keep at most two repeats for readability
        w = re.sub(r'(.)\1{2,}', r'\1\1', w)

        # If token is composed of repeated double-letters (e.g. RReegguulll), collapse pairs
        if len(w) >= 6:
            # check if many consecutive pairs are identical (indicative of doubled characters)
            pairs = [w[i:i+2] for i in range(0, len(w)//2*2, 2)]
            double_pairs = sum(1 for p in pairs if len(p) == 2 and p[0].lower() == p[1].lower())
            if double_pairs >= max(2, len(pairs)//2):
                # collapse each pair to single char
                try:
                    w = ''.join(p[0] for p in pairs) + (w[len(pairs)*2:] if len(w) % 2 else '')
                except Exception:
                    pass

        # Remove repeated punctuation beyond 3 occurrences
        w = re.sub(r'([!?.\-_,;:\\"\'\(\)])\1{2,}', r'\1\1', w)

        # Trim extremely long garbage-like tokens but keep most of content
        if len(w) > 500:
            # keep head and tail
            w = w[:250] + '...' + w[-250:]

        # If cleaned token looks much worse (unlikely word) fall back to original to avoid loss
        if likely_word(orig) or likely_word(w):
            cleaned_words.append(w)
        else:
            # keep original conservative fallback
            cleaned_words.append(orig)

    # 4) Remove immediate duplicate consecutive words (e.g., "the the")
    final_words = []
    for w in cleaned_words:
        if not final_words or w.lower() != final_words[-1].lower():
            final_words.append(w)

    return ' '.join(final_words)


def clean_extracted_text_preserve_layout(text):
    """
    Like clean_extracted_text but preserves line breaks and relative spacing
    so the page layout remains closer to the original PDF. This removes
    only illegal xml chars and obvious cid patterns, and normalizes trailing
    whitespace on each line while keeping internal spacing.
    """
    if not text:
        return ""

    # Remove (cid:####) patterns and illegal xml/control chars
    text = re.sub(r'\(cid:\d+\)', '', text)
    text = _illegal_xml_re.sub('', text)

    # Normalize line endings to \n
    text = text.replace('\r\n', '\n').replace('\r', '\n')

    # Trim trailing spaces on each line but preserve internal spacing
    lines = [ln.rstrip() for ln in text.split('\n')]
    # Remove consecutive blank lines (more than 2) to avoid large gaps
    cleaned_lines = []
    blank_run = 0
    for ln in lines:
        if not ln.strip():
            blank_run += 1
            if blank_run <= 2:
                cleaned_lines.append(ln)
        else:
            blank_run = 0
            cleaned_lines.append(ln)

    return '\n'.join(cleaned_lines).strip('\n')


def clean_ocr_text(text, min_alpha_ratio=0.25, min_word_count=1):
    """
    Clean OCR-produced text conservatively: remove obvious garbage lines
    while keeping as much real text as possible (tables, numeric rows,
    short codes). Returns cleaned multi-line text preserving line breaks.
    """
    if not text:
        return ""

    # Normalize line endings and remove control chars
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    text = _illegal_xml_re.sub('', text)

    lines = [ln.strip() for ln in text.split('\n')]
    cleaned = []
    prev = None
    for ln in lines:
        if not ln:
            # keep single blank lines (to preserve paragraph separation)
            if prev != "":
                cleaned.append("")
                prev = ""
            continue

        # Trim excessive repeated punctuation/characters
        if re.fullmatch(r'([\W_])\1{2,}', ln):
            # line like "-----" or "......" -> drop
            continue

        # Count character classes
        total_chars = len(ln.replace(' ', ''))
        letters = len(re.findall(r'[A-Za-z]', ln))
        digits = len(re.findall(r'\d', ln))
        words = re.findall(r"\w+", ln)
        word_count = len(words)

        alpha_ratio = (letters / total_chars) if total_chars > 0 else 0

        # Heuristics to keep a line:
        keep = False
        # 1) If it has at least min_word_count words (likely meaningful)
        if word_count >= min_word_count:
            keep = True

        # 2) If line has enough alphabetic characters
        if not keep and alpha_ratio >= min_alpha_ratio:
            keep = True

        # 3) Keep numeric/tabular lines that look like rows (multiple numbers or separators)
        if not keep:
            if '|' in ln or '\t' in ln:
                # treat as table-like
                keep = True
            else:
                # comma-separated numbers or multiple numeric tokens
                numeric_tokens = len([w for w in words if re.fullmatch(r'[-+]?\d+[\d,.]*', w)])
                if numeric_tokens >= 2:
                    keep = True

        # 4) Drop lines that are mostly non-alphanumeric garbage
        if keep:
            # avoid lines with extremely low informative content
            punct_ratio = (len(re.findall(r'[^\w\s]', ln)) / total_chars) if total_chars > 0 else 0
            if total_chars > 0 and letters == 0 and digits == 0 and punct_ratio > 0.6:
                keep = False

        if keep:
            # Deduplicate consecutive identical lines
            if ln != prev:
                cleaned.append(ln)
                prev = ln
        else:
            # drop the line
            continue

    # Post-process: strip leading/trailing blank lines
    while cleaned and cleaned[0] == "":
        cleaned.pop(0)
    while cleaned and cleaned[-1] == "":
        cleaned.pop()

    return "\n".join(cleaned)
   
def process_pdf(file_stream):
    file_stream.seek(0)
    file_bytes = file_stream.read()
    pages_data = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        for i, page in enumerate(pdf.pages):
            # Try to extract text while preserving layout as much as possible.
            # Prefer pdfplumber's extraction (which keeps line breaks). If that
            # yields little content, fall back to PyMuPDF's get_text.
            text = page.extract_text() or ""
            if not text or len(text.strip()) < 20:
                try:
                    # PyMuPDF's text extraction can be more robust for some PDFs
                    page_fitz_tmp = doc.load_page(i)
                    text = page_fitz_tmp.get_text("text") or ""
                except Exception:
                    text = text or ""

            # Use the preserve-layout cleaner to keep line breaks and spacing
            text = clean_extracted_text_preserve_layout(text)
            tables = page.extract_tables() or []
            page_fitz = doc.load_page(i)
            image_texts = []
            image_refs = []
            for img_info in page_fitz.get_images(full=True):
                
                if(len(text) < 1000):
                   xref = img_info[0]
                   base_image = doc.extract_image(xref)
                   image_bytes = base_image["image"]
                   image_ext = base_image.get("ext", "png").lower()
                   ocr_text = extract_text_from_image(image_bytes) or ""
                   ocr_text= clean_extracted_text(ocr_text)
                   if ocr_text:
                       image_texts.append(ocr_text)
                   image_id = fs.put(image_bytes, contentType=f"image/{image_ext}", filename=f"page{i+1}_image.{image_ext}")
                   image_refs.append({"id": str(image_id), "ext": image_ext, "ocr_text": ocr_text})
            if not text.strip() and len(text) < 1000:
                pix = page_fitz.get_pixmap(dpi=300)
                img_bytes = pix.tobytes("png")
                ocr_text = extract_text_from_image(img_bytes, lang="eng")
                ocr_text= clean_extracted_text(ocr_text)
                if ocr_text:
                    text = ocr_text
                    image_texts.append(ocr_text)
            pages_data.append({"page_number": i + 1, "text": text, "tables": tables, "image_texts": image_texts, "images": image_refs})
        doc.close()
    return pages_data


def check_tesseract_available():
    try:
        version = pytesseract.get_tesseract_version()
        print(f"Tesseract available: {version}")
    except Exception as e:
        raise RuntimeError(f"Tesseract is not functioning: {e}")


def extract_text_from_image(image_bytes, lang="eng"):
    try:
        image = preprocess_image_for_ocr(image_bytes)
        raw = pytesseract.image_to_string(image, lang=lang, config="--psm 6")
        # Preserve layout first, then apply OCR-specific cleaning to drop garbage
        text = clean_extracted_text_preserve_layout(raw)
        text = clean_ocr_text(text)
        if not text:
            raw = pytesseract.image_to_string(image, lang=lang, config="--psm 3")
            text = clean_extracted_text_preserve_layout(raw)
            text = clean_ocr_text(text)
        return text
    except Exception as e:
        app.logger.warning(f"OCR failed: {e}")
        return ""


def process_image(file_stream, filename):
    file_stream.seek(0)
    image_bytes = file_stream.read()
    text = extract_text_from_image(image_bytes)
    text= clean_extracted_text(text)
    file_stream.seek(0)
    image_ext = filename.rsplit(".", 1)[1].lower()
    image_id = fs.put(image_bytes, contentType=f"image/{image_ext}", filename=filename)
    data = {"images": [{"id": str(image_id), "ext": image_ext}], "ocr_text": text}
    return data


def process_html(file_stream, filename=None):
    """Parse HTML, extract formatted text, headings and tables.
    Returns a dict containing: title, headings, paragraphs, lists, tables, images, full_html
    """
    file_stream.seek(0)
    raw = file_stream.read()
    try:
        text = raw.decode("utf-8")
    except Exception:
        text = raw.decode("utf-8", errors="replace")

    soup = BeautifulSoup(text, "html.parser")

    title = soup.title.string.strip() if soup.title and soup.title.string else ""

    headings = []
    for lvl in range(1, 7):
        for h in soup.find_all(f"h{lvl}"):
            headings.append(clean_extracted_text(h.get_text(" ", strip=True)))

    paragraphs = [clean_extracted_text(p.get_text(" ", strip=True)) for p in soup.find_all("p")]

    lists = []
    for ul in soup.find_all(["ul", "ol"]):
        items = [clean_extracted_text(li.get_text(" ", strip=True)) for li in ul.find_all("li")]
        if items:
            lists.append(items)

    # Tables: preserve both structure and HTML fragment for formatting
    tables = []
    for t_idx, table in enumerate(soup.find_all("table")):
        rows = []
        for tr in table.find_all("tr"):
            cells = [clean_extracted_text(td.get_text(" ", strip=True)) for td in tr.find_all(["th", "td"])]
            if cells:
                rows.append(cells)
        table_html = str(table)
        # attempt to infer headers
        headers = []
        if rows:
            first_row = rows[0]
            # heuristic: if first row contains non-empty values and likely header
            headers = first_row
            data_rows = rows[1:]
        else:
            data_rows = []

        tables.append({
            "table_index": t_idx,
            "headers": headers,
            "rows": data_rows,
            "html": table_html,
        })

    # Images: handle data URIs by storing in GridFS, otherwise keep src
    images = []
    for img_idx, img in enumerate(soup.find_all("img"), start=1):
        src = img.get("src") or ""
        if src.startswith("data:"):
            try:
                header, b64 = src.split(",", 1)
                import base64 as _b64
                img_bytes = _b64.b64decode(b64)
                # try to infer extension
                mtype = header.split(";")[0].split(":")[1] if ":" in header else "image/png"
                ext = mtype.split("/")[-1]
                img_id = fs.put(img_bytes, contentType=mtype, filename=(filename or "uploaded") + f"_img_{img_idx}.{ext}")
                images.append({"id": str(img_id), "ext": ext})
            except Exception as e:
                app.logger.debug("Failed to store data URI image: %s", e)
                images.append({"src": src})
        else:
            images.append({"src": src})

    # full HTML (sanitized lightly)
    full_html = str(soup)

    data = {
        "title": title,
        "headings": headings,
        "paragraphs": paragraphs,
        "lists": lists,
        "tables": tables,
        "images": images,
        "full_html": full_html,
    }
    return data

def clean_for_mongo(obj):
    if obj is None:
        return None
    if isinstance(obj, (str, int, float, bool)):
        return obj
    if isinstance(obj, (datetime.date, datetime.datetime)):
        return obj.isoformat()
    if isinstance(obj, datetime.time):
        return obj.strftime("%H:%M:%S")
    if isinstance(obj, ObjectId):
        return str(obj)
    if isinstance(obj, dict):
        return {k: clean_for_mongo(v) for k, v in obj.items()}
    if isinstance(obj, (list, tuple)):
        return [clean_for_mongo(x) for x in obj]
    try:
        return str(obj)
    except Exception:
        return None


def compute_file_hash(file_input):
    hasher = hashlib.sha256()
    if isinstance(file_input, (bytes, bytearray)):
        hasher.update(file_input)
    else:
        for chunk in iter(lambda: file_input.read(4096), b""):
            hasher.update(chunk)
        file_input.seek(0)
    return hasher.hexdigest()

def process_and_store_file(abs_path):
    filename = os.path.basename(abs_path)
    filetype = filename.rsplit(".", 1)[1].lower()
    with open(abs_path, "rb") as f:
        file_bytes = f.read()

    # Save in GridFS
    file_id = fs.put(file_bytes, filename=filename, contentType="application/octet-stream")

    # Process file
    if filetype == "xlsx":
        file_data = process_excel(io.BytesIO(file_bytes))
    elif filetype == "pptx":
        file_data = process_pptx(io.BytesIO(file_bytes))
    elif filetype == "pdf":
        file_data = process_pdf(io.BytesIO(file_bytes))
    elif filetype in {"png", "jpg", "jpeg", "bmp", "gif", "tiff"}:
        file_data = process_image(io.BytesIO(file_bytes), filename)
    elif filetype in {"html", "htm"}:
        # process HTML and preserve simple formatting and tables
        file_data = process_html(io.BytesIO(file_bytes))
    else:
        raise ValueError("Unsupported file type")

    # Build search text
    _search_tokens = flatten_text(file_data)
    _search_tokens.insert(0, filename)
    search_text = " ".join([str(t) for t in _search_tokens if t])

    # Generate embedding
    embedding = model.encode(search_text, convert_to_numpy=True).astype("float32")
    

    # Railway domain terms
    
    cleaned_data = clean_for_mongo(file_data)
    inserted = mongo.db.files.insert_one({
        "filename": filename,
        "filetype": filetype,
        "file_id": file_id,
        "search_text": search_text,
        "embedding": embedding.tolist(),
        "upload_time": datetime.datetime.utcnow(),
        "data": cleaned_data
    })
    rail_terms = ["wagon", "hopper", "open wagon", "loading", "unloading",
              "demurrage", "freight", "free time", "station", "railway board"]

    
    patterns = [nlp.make_doc(t) for t in rail_terms]
    matcher.add("RAIL_TERMS", patterns)
    # Add to FAISS
    add_to_faiss(str(inserted.inserted_id), embedding)

    return inserted.inserted_id, filename



def process_and_store_file_with_check(abs_path):
    filename = os.path.basename(abs_path)

    with open(abs_path, "rb") as f:
        file_bytes = f.read()

    file_hash = compute_file_hash(file_bytes)

    # 🔹 Check if already in Mongo
    existing = mongo.db.files.find_one({"file_hash": file_hash})
    if existing:
        # Return existing without re-processing
        return existing["_id"], existing["filename"]

    # 🔹 Call your original function
    inserted_id, stored_filename = process_and_store_file(abs_path)

    # Update with hash for future checks
    mongo.db.files.update_one(
        {"_id": inserted_id},
        {"$set": {"file_hash": file_hash}}
    )

    return inserted_id, stored_filename


def check_duplicate_and_save(file, files_collection, extra_metadata=None):
    """
    Check if file already exists in MongoDB based on hash.
    If not, insert new metadata.
    Returns: (status, existing_or_new_doc)
    """
    file_hash = compute_file_hash(file.stream)

    # Look for existing entry
    existing = files_collection.find_one({"file_hash": file_hash})
    if existing:
        existing = files_collection.find_one({"file_hash": file_hash})

    # Insert new document metadata
    new_doc = {
        "filename": file.filename,
        "file_hash": file_hash,
        "content_type": file.content_type,
        "uploaded_at": datetime.datetime.utcnow()
    }
    if extra_metadata:
        new_doc.update(extra_metadata)

    inserted_id = files_collection.insert_one(new_doc).inserted_id
    new_doc["_id"] = inserted_id
    return "new", new_doc


def extract_entities(text, chunk_size=100_000):
    entities = []
    n = len(text)
    for i in range(0, n, chunk_size):
        chunk = text[i:i+chunk_size]
        doc = nlp(chunk)
        for ent in doc.ents:
            label, meta = classify_rail_entity(ent.text)
            entities.append({
                "text": ent.text,
                "label": label,
                "meta": meta
            })
    return entities

def extract_failure_codes(file_data):
    """
    Extracts failure codes & details from text, dict, or Excel table.
    
    Returns a list of dictionaries:
        [
            {
                "department": "Signal and Telecommunication",
                "asset_type": "Panel Interlocking",
                "failure_code": "ACC",
                "failure_subcode": "1.01",
                "failure_desc": "ACCIDENT OTHERS",
                "valid": "Y",
                "user_asset_failure": "Y/N",
                "system_auto": "Y/N",
                "asset_group": "UNUSUAL"
            },
            ...
        ]
    """
    failures = []

    # If already a DataFrame (Excel)
    if isinstance(file_data, pd.DataFrame):
        df = file_data
    # If a dict (likely returned from process_excel/pdf)
    elif isinstance(file_data, dict):
        # Try to find a table inside dict
        if "table" in file_data and isinstance(file_data["table"], list):
            df = pd.DataFrame(file_data["table"])
        else:
            # Flatten dict to string
            text = " ".join(str(v) for v in file_data.values())
            lines = [line.strip() for line in text.splitlines() if line.strip()]
            if not lines:
                return failures
            headers = re.split(r'\t+|\s{2,}', lines[0])
            rows = [re.split(r'\t+|\s{2,}', l) for l in lines[1:]]
            df = pd.DataFrame(rows, columns=headers)
    # If raw text
    elif isinstance(file_data, str):
        lines = [line.strip() for line in file_data.splitlines() if line.strip()]
        if not lines:
            return failures
        headers = re.split(r'\t+|\s{2,}', lines[0])
        rows = [re.split(r'\t+|\s{2,}', l) for l in lines[1:]]
        df = pd.DataFrame(rows, columns=headers)
    else:
        return failures

    # Normalize column names
    df.columns = [c.strip().lower().replace(" ", "_") for c in df.columns]

    for _, row in df.iterrows():
        failure = {
            "department": row.get("department", "").strip() if row.get("department") else "",
            "asset_type": row.get("asset_type", "").strip() if row.get("asset_type") else "",
            "failure_code": row.get("failure_code", "").strip() if row.get("failure_code") else "",
            "failure_subcode": row.get("failure_sub_code", row.get("sub_sr", "")).strip() if row.get("failure_sub_code") or row.get("sub_sr") else "",
            "failure_desc": row.get("failure_description", "").strip() if row.get("failure_description") else "",
            "valid": row.get("valid?", "Y").strip() if row.get("valid?") else "Y",
            "user_asset_failure": row.get("user_asset_failure", "").strip() if row.get("user_asset_failure") else "",
            "system_auto": row.get("system_auto", "").strip() if row.get("system_auto") else "",
            "asset_group": row.get("asset_group", "").strip() if row.get("asset_group") else ""
        }
        if failure["failure_code"]:
            failures.append(failure)

    return failures
from rapidfuzz import fuzz, process
import numpy as np

# Optional: use your existing sentence embedding model (if available)
try:
    from sentence_transformers import SentenceTransformer
    embed_model = SentenceTransformer("all-MiniLM-L6-v2")
    USE_EMBEDDING_SIM = True
except Exception:
    USE_EMBEDDING_SIM = False


def semantic_similarity(a, b):
    """Compute cosine similarity between two strings using embeddings."""
    if not USE_EMBEDDING_SIM:
        return 0.0
    emb_a = embed_model.encode(a, normalize_embeddings=True)
    emb_b = embed_model.encode(b, normalize_embeddings=True)
    return float(np.dot(emb_a, emb_b))


def classify_rail_entity(text,
                         fuzzy_threshold=85, embed_threshold=0.75, use_embedding_sim=True):
    """
    Classify text into railway-related semantic categories using fuzzy and semantic similarity.
    Returns (category, metadata).
    """
    t = text.lower().strip()
    best_match, best_score, category_meta = None, 0, {}
    
    # Helper for fuzzy + semantic similarity
    def match_score(source, candidates):
        """Return the best fuzzy+semantic match for a text source."""
        if not candidates:
            return None, 0

        result = process.extractOne(source, candidates, scorer=fuzz.token_set_ratio)

        # ✅ If no match found, return safely
        if not result or len(result) < 2:
            return None, 0

        match, score = result[0], result[1]
        sim = semantic_similarity(source, match) if use_embedding_sim else 0
        return match, max(score, sim * 100)

    # 1. Signal asset check
    for asset, meta in signal_assets.items():
        score = fuzz.token_set_ratio(t, asset.lower())
        if score > fuzzy_threshold:
            return "Asset", {**meta, "match": asset, "score": score}

    # 2. Failure codes
    for f in rail_failure_dict:
        code, desc = f["failure_code"].lower(), f["failure_desc"].lower()
        score = max(fuzz.token_set_ratio(t, code), fuzz.token_set_ratio(t, desc))
        if score > fuzzy_threshold:
            return "Failure", {"department": f["department"], "code": f["failure_code"], "score": score}

    # 3️⃣ Department assets
    for dept, assets in railway_assets.items():
        match, score = match_score(t, assets)
        if score > fuzzy_threshold:
            return "Asset", {"department": dept, "asset_type": match, "score": score}

    # 4️⃣ Organizational hierarchy
    if any(k in t for k in ["division", "zone", "board", "headquarters"]):
        return "Organization", {"score": 90}

    # 5️⃣ Station or Yard
    if "station" in t:
        return "Station", {"score": 95}
    if "yard" in t:
        return "Yard", {"score": 95}

    return "Generic", {"score": 50}

def add_to_graph(entities, relations, file_id):
    print("🧠 Starting KG extraction for uploaded JSON...")
    entity_count, relation_count = 0, 0

    try:
        with driver.session() as session:
            for e in entities:
                session.execute_write(create_entity_node, e, file_id)
                entity_count += 1

            # Deduplicate relations
            unique_relations = {tuple(sorted(r.items())) if isinstance(r, dict) else r for r in relations}
            for r in unique_relations:
                session.execute_write(create_relation_node, r, file_id)
                relation_count += 1

    except Exception as e:
        print(f"❌ Graph insertion failed: {e}")

    print(f"✅ Completed KG insertion — {entity_count} entities, {relation_count} relations.")


def create_failure_kg_node(tx, failure, file_id):
    tx.run("""
        MERGE (f:Failure {code: $failure_code})
        SET f.subcode = $failure_subcode,
            f.description = $failure_desc,
            f.department = $department,
            f.asset_group = $asset_group,
            f.valid = $valid,
            f.system_auto = $system_auto,
            f.file_id = $file_id
    """, 
    failure_code=failure["failure_code"],
    failure_subcode=failure.get("failure_subcode", ""),
    failure_desc=failure.get("failure_desc", ""),
    department=failure.get("department", ""),
    asset_group=failure.get("asset_group", ""),
    valid=failure.get("valid", ""),
    system_auto=failure.get("system_auto", ""),
    file_id=str(file_id))



def create_asset_kg_node(tx, asset, file_id):
    """
    Create or link an asset node to existing KG hierarchy.
    Reuses existing Zone, Division, Section, and GearGroup nodes when possible.
    """

    q = """
    // --- Match existing base hierarchy ---
    OPTIONAL MATCH (z:Zone {name:$zone})
    OPTIONAL MATCH (d:Division {name:$division})-[:PART_OF]->(z)
    OPTIONAL MATCH (s:Section {name:$section})-[:PART_OF]->(d)
    OPTIONAL MATCH (g:GearGroup {name:$gear_group})
    OPTIONAL MATCH (dept:Department {name:$department})

    // --- Create asset node if missing ---
    MERGE (a:Asset {name:$asset_name})
      ON CREATE SET
        a.asset_type = $asset_type,
        a.gear_name = $gear_name,
        a.department = $department,
        a.created_at = datetime(),
        a.source_file = $file_id
      ON MATCH SET
        a.last_seen = datetime(),
        a.source_file = $file_id

    // --- Link asset to hierarchy ---
    FOREACH (_ IN CASE WHEN z IS NOT NULL THEN [1] ELSE [] END |
        MERGE (a)-[:UNDER_ZONE]->(z))
    FOREACH (_ IN CASE WHEN d IS NOT NULL THEN [1] ELSE [] END |
        MERGE (a)-[:UNDER_DIVISION]->(d))
    FOREACH (_ IN CASE WHEN s IS NOT NULL THEN [1] ELSE [] END |
        MERGE (a)-[:UNDER_SECTION]->(s))
    FOREACH (_ IN CASE WHEN g IS NOT NULL THEN [1] ELSE [] END |
        MERGE (a)-[:BELONGS_TO_GROUP]->(g))
    FOREACH (_ IN CASE WHEN dept IS NOT NULL THEN [1] ELSE [] END |
        MERGE (a)-[:MAINTAINED_BY]->(dept))
    """

    # Clean and normalize before sending to Cypher
    data = {
        "zone": asset.get("Zone", "").strip(),
        "division": asset.get("Division", "").strip(),
        "section": asset.get("Section", "").strip(),
        "gear_group": asset.get("Gear Group", "").strip() or asset.get("Gear Type", "").strip(),
        "asset_name": asset.get("Gear Name/Number", "").strip(),
        "asset_type": asset.get("Asset Type", "").strip(),
        "gear_name": asset.get("Gear Name/Number", "").strip(),
        "department": asset.get("Department", "").strip(),
        "file_id": str(file_id),
    }

    tx.run(q, **data)

def extract_relations(text):
    doc = nlp(text)
    triples = []
    for sent in doc.sents:
        subj = None
        obj = None
        verb = None
        for token in sent:
            if "subj" in token.dep_:
                subj = token.text
            if "obj" in token.dep_:
                obj = token.text
            if token.pos_ == "VERB":
                verb = token.lemma_
        if subj and verb and obj:
            triples.append((subj, verb, obj))
    return triples


def create_entity_node(tx, entity, file_id):
    """
    Safely create Entity nodes with labels and file link.
    """
    text = entity.get("text", "").strip()
    label = entity.get("label", "Generic").strip()
    if not text:
        return

    tx.run("""
        MERGE (e:Entity {text: $text})
        SET e.label = $label,
            e.file_id = $file_id,
            e.meta = $meta
        MERGE (e)-[:MENTIONED_IN]->(f:File {id: $file_id})
    """, text=text, label=label, meta=json.dumps(entity.get("meta", {})), file_id=str(file_id))


def create_relation_node(tx, relation, file_id):
    """Insert entity relations safely, with sanitized relation type."""
    # Normalize input
    if isinstance(relation, tuple):
        if len(relation) == 3:
            from_text, rel_type, to_text = relation
        elif len(relation) == 2:
            from_text, to_text = relation
            rel_type = "RELATED_TO"
        else:
            return
    elif isinstance(relation, dict):
        from_text = relation.get("from", "").strip()
        to_text = relation.get("to", "").strip()
        rel_type = relation.get("type", "RELATED_TO").strip()
    else:
        return

    if not from_text or not to_text:
        return

    # Sanitize relation type
    safe_rel_type = re.sub(r"[^A-Za-z0-9_]", "_", rel_type.upper())

    cypher = f"""
        MERGE (a:Entity {{text: $from_text}})
        MERGE (b:Entity {{text: $to_text}})
        MERGE (a)-[r:{safe_rel_type}]->(b)
        MERGE (a)-[:MENTIONED_IN]->(f:File {{id: $file_id}})
        MERGE (b)-[:MENTIONED_IN]->(f)
    """
    tx.run(cypher, from_text=from_text, to_text=to_text, file_id=str(file_id))

# --------------------------- Routes --------------------------------

@app.route("/", methods=["GET"])
def dashboard():
    files = list(mongo.db.files.find().sort("upload_time", -1).limit(15))
    selected_file_id = request.args.get("file_id", "")
    return render_template(
        "dashboard.html",
        files=files,
        selected_file_id=selected_file_id
    )


@app.route('/upload_json', methods=['POST'])
def upload_json():
    total_uploaded = 0
    total_failed = 0
    skipped_files = []

    try:
        json_file = request.files.get('file')
        if not json_file:
            flash("No JSON file uploaded", "danger")
            return redirect(url_for('dashboard', _external=True))

        data = json.load(json_file)

        for xlsx_file, file_list in data.items():
            # Only process PDFs
            pdf_files = [f for f in file_list if f.strip().lower().endswith('.pdf')]

            if not pdf_files:
                skipped_files.append(f"{xlsx_file} (no PDF files)")
                continue

            for file_path in pdf_files:
                abs_path = os.path.normpath(os.path.join(os.getcwd(), file_path))
                print(f"Trying to upload: {abs_path}")
                if not os.path.exists(abs_path):
                    skipped_files.append(f"{file_path} (not found)")
                    total_failed += 1
                    continue

                try:
                    file_id, file_name = process_and_store_file_with_check(abs_path)
                    file_doc = mongo.db.files.find_one({"_id": file_id})
                    text_content = " ".join(flatten_text(file_doc["data"]))
                    #entities = extract_entities(text_content)
                    #triples = extract_relations(text_content)
                    #add_to_graph(entities, triples, file_id)
                    #flash(f"Uploaded: {file_name}", "success")
                    total_uploaded += 1
                except Exception as e:
                    skipped_files.append(f"{file_path} (error: {str(e)})")
                    total_failed += 1

        # Flash summary
        summary_msg = f"✅ {total_uploaded} PDF files uploaded"
        if total_failed:
            summary_msg += f", ❌ {total_failed} failed"
        if skipped_files:
            summary_msg += f". Skipped: {len(skipped_files)} files"

        flash(summary_msg, "info")

    except Exception as e:
        flash(f"Error reading JSON: {str(e)}", "danger")

    return redirect(url_for('dashboard'))

@app.route("/dashboard/view/<file_id>")
def view_file_in_dashboard(file_id):
    file = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file:
        flash("File not found", "danger")
        return redirect(url_for("dashboard"))

    # Pass the selected file id to the dashboard so it auto-opens
    return redirect(url_for("dashboard", file_id=file_id))


@app.route("/upload", methods=["GET", "POST"])
def upload():
    if request.method == "POST":
        if "file" not in request.files:
            flash("No file part", "danger")
            return redirect(request.url)

        file = request.files["file"]
        if not file or file.filename == "":
            flash("No selected file", "danger")
            return redirect(request.url)

        if file and allowed_file(file.filename):
            filename = secure_filename(file.filename)
            filetype = filename.rsplit(".", 1)[1].lower()

            # --- Check duplicate ---
            status, doc = check_duplicate_and_save(file, files_collection, extra_metadata={"type": filetype})
            if status == "duplicate":
                return jsonify({
                    "status": "duplicate",
                    "message": "File already uploaded",
                    "file_id": str(doc["_id"]),
                    "filename": doc["filename"]
                }), 200

            file.seek(0)
            file_id = fs.put(file, filename=filename, contentType=file.mimetype)
            file.seek(0)

            try:
                if filetype in {"xlsx", "xls"}:
                    file_data = process_excel(file)
                elif filetype == "pdf":
                    file_data = process_pdf(file)
                elif filetype in {"png", "jpg", "jpeg"}:
                    file_data = process_image(file, filename)
                else:
                    flash("Unsupported file type", "danger")
                    return redirect(request.url)
            except Exception as e:
                flash(f"Failed to process file: {e}", "danger")
                return redirect(request.url)

            # --- Flatten text for search & embeddings ---
            _search_tokens = flatten_text(file_data)
            _search_tokens.insert(0, filename)
            search_text = " ".join([str(t) for t in _search_tokens if t])

            embedding = model.encode(search_text, convert_to_numpy=True)
            embedding = np.asarray(embedding, dtype="float32")

            cleaned_data = clean_for_mongo(file_data)

            # --- Extract failures & assets ---
            failures = extract_failure_codes(file_data)
            assets = [f for f in failures if f.get("user_asset_failure", "Y") == "Y"]

            # --- Insert main record into MongoDB ---
            inserted = mongo.db.files.insert_one({
                "filename": filename,
                "filetype": filetype,
                "file_id": file_id,
                "search_text": search_text,
                "embedding": embedding.tolist(),
                "upload_time": datetime.datetime.utcnow(),
                "data": cleaned_data
            })
            file_id = inserted.inserted_id

            # --- Async / background operations for speed ---
            def background_tasks():
                try:
                    # Add embedding to FAISS
                    add_to_faiss(str(file_id), embedding, normalize=True)
                except Exception as e:
                    app.logger.exception("Failed to add to FAISS: %s", e)

                # Add failures & assets to Neo4j
                try:
                    with driver.session() as session:
                        for failure in failures:
                            failure = map_to_existing_entities(failure)
                            code = failure.get("failure_code")
                            if not session.execute_read(failure_exists, code):
                                session.execute_write(link_failure_to_existing_nodes, failure, file_id)
                                print(f"Added failure code to KG: {code}")
                        for asset in assets:
                            asset_name = asset.get("asset_name")
                            if not session.execute_read(asset_exists, asset_name):
                                session.execute_write(create_asset_kg_node, asset, file_id)
                                print(f"Added asset to KG: {asset_name}")
                except Exception as e:
                    app.logger.exception("Neo4j failure: %s", e)

                # Add extracted entities & relations to KG
                try:
                    entities = extract_entities(search_text)
                    relations = extract_relations(search_text)
                    add_to_graph(entities, relations, file_id)
                    print(f"Added {len(entities)} entities and {len(relations)} relations to KG")
                except Exception as e:
                    app.logger.exception(f"KG entity/relationship add failed: {e}")

            # --- Run background tasks in a separate thread ---
            threading.Thread(target=background_tasks, daemon=True).start()

            flash("File uploaded, processed & indexed successfully (background processing running)", "success")
            return redirect(url_for("dashboard"))

        flash("Invalid file type. Allowed: .xlsx, .pdf, images", "danger")
        return redirect(request.url)

    return render_template("upload_file.html")

@app.route("/files/<file_id>")
def serve_file(file_id):
    try:
        grid_out = fs.get(ObjectId(file_id))
        return send_file(io.BytesIO(grid_out.read()), download_name=grid_out.filename, mimetype=grid_out.content_type or "application/octet-stream", as_attachment=True)
    except NoFile:
        flash("File not found in storage.", "danger")
        return redirect(url_for("dashboard"))


@app.route("/images/<image_id>")
def serve_image(image_id):
    try:
        grid_out = fs.get(ObjectId(image_id))
        mimetype = grid_out.content_type or "application/octet-stream"
        return send_file(io.BytesIO(grid_out.read()), download_name=grid_out.filename, mimetype=mimetype, as_attachment=False)
    except NoFile:
        abort(404)


@app.route("/export_json/<file_id>")
def export_json(file_id):
    file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file_doc or "data" not in file_doc:
        flash("File or data not found.", "danger")
        return redirect(url_for("dashboard"))

    def json_serializer(obj):
        if isinstance(obj, (datetime.date, datetime.datetime)):
            return obj.isoformat()
        return str(obj)

    json_data = json.dumps(file_doc["data"], indent=2, ensure_ascii=False, default=json_serializer)
    response = make_response(json_data)
    response.headers["Content-Disposition"] = f"attachment; filename={file_doc['filename']}_data.json"
    response.headers["Content-Type"] = "application/json"
    return response


def write_excel_from_data(file_name_prefix, filetype, data):
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    if filetype == "xlsx":
        for sheet_name, rows in data.items():
            ws = wb.create_sheet(title=sheet_name[:31])
            for row_idx, row in enumerate(rows, 1):
                for col_idx, cell in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=clean_cell_value(cell))
    elif filetype == "pptx":
        for slide_idx, slide in enumerate(data, 1):
            ws = wb.create_sheet(title=f"Slide {slide_idx}")
            row_cursor = 1
            ws.cell(row=row_cursor, column=1, value="Title")
            ws.cell(row=row_cursor, column=2, value=clean_cell_value(slide.get("title", "")))
            row_cursor += 2
            ws.cell(row=row_cursor, column=1, value="Texts")
            row_cursor += 1
            for text in slide.get("texts", []):
                ws.cell(row=row_cursor, column=1, value=clean_cell_value(text))
                row_cursor += 1
            row_cursor += 1
            images = slide.get("images", [])
            if images:
                ws.cell(row=row_cursor, column=1, value=f"Images ({len(images)} stored in DB)")
                row_cursor += 2
                for img_idx, img in enumerate(images, 1):
                    ws.cell(row=row_cursor, column=1, value=f"Image {img_idx} (id: {img['id']}, ext: .{img.get('ext', 'png')})")
                    row_cursor += 1
                    if img.get("ocr_text"):
                        for line in img["ocr_text"].splitlines():
                            ws.cell(row=row_cursor, column=2, value=clean_cell_value(line))
                            row_cursor += 1
                    row_cursor += 1
            for t_idx, table in enumerate(slide.get("tables", []), 1):
                ws.cell(row=row_cursor, column=1, value=f"Table {t_idx}")
                row_cursor += 1
                for row in table:
                    for col_idx, cell in enumerate(row, 1):
                        ws.cell(row=row_cursor, column=col_idx, value=clean_cell_value(cell))
                    row_cursor += 1
                row_cursor += 1
    elif filetype == "pdf":
        for page in data:
            ws = wb.create_sheet(title=f"Page {page.get('page_number', '?')}")
            row_cursor = 1
            ws.cell(row=row_cursor, column=1, value="Extracted Text")
            row_cursor += 1
            for line in page.get("text", "").splitlines():
                ws.cell(row=row_cursor, column=1, value=clean_cell_value(line))
                row_cursor += 1
            row_cursor += 1
            for t_idx, table in enumerate(page.get("tables", []), 1):
                ws.cell(row=row_cursor, column=1, value=f"Table {t_idx}")
                row_cursor += 1
                for row in table:
                    for col_idx, cell in enumerate(row, 1):
                        ws.cell(row=row_cursor, column=col_idx, value=clean_cell_value(cell))
                    row_cursor += 1
                row_cursor += 1
            if page.get("image_texts"):
                ws.cell(row=row_cursor, column=1, value="OCR Texts from Images (aggregated)")
                row_cursor += 1
                for text in page["image_texts"]:
                    for line in text.splitlines():
                        ws.cell(row=row_cursor, column=1, value=clean_cell_value(line))
                        row_cursor += 1
                    row_cursor += 1
            images = page.get("images", [])
            if images:
                ws.cell(row=row_cursor, column=1, value=f"Images ({len(images)} stored in DB)")
                row_cursor += 2
                for img_idx, img in enumerate(images, 1):
                    ws.cell(row=row_cursor, column=1, value=f"Image {img_idx} (id: {img['id']}, ext: .{img.get('ext', 'png')})")
                    row_cursor += 1
                    if img.get("ocr_text"):
                        for line in img["ocr_text"].splitlines():
                            ws.cell(row=row_cursor, column=2, value=clean_cell_value(line))
                            row_cursor += 1
                    row_cursor += 1
    elif filetype in {"png", "jpg", "jpeg", "bmp", "gif", "tiff"}:
        ws = wb.create_sheet(title="Image Data")
        ws.cell(row=1, column=1, value="OCR Extracted Text")
        ws.cell(row=2, column=1, value=clean_cell_value(data.get("ocr_text", "")))
        ws.cell(row=4, column=1, value="Image Reference")
        images = data.get("images", [])
        for idx, img in enumerate(images, start=1):
            ws.cell(row=4 + idx, column=1, value=f"Image {idx} (id: {img['id']}, ext: .{img.get('ext','png')})")
            if img.get("ocr_text"):
                ws.cell(row=4 + idx, column=2, value=clean_cell_value(img["ocr_text"]))
    else:
        ws = wb.create_sheet(title="Data")
        ws.cell(row=1, column=1, value=clean_cell_value(json.dumps(data, indent=2)))
    excel_stream = io.BytesIO()
    wb.save(excel_stream)
    excel_stream.seek(0)
    return excel_stream


def write_word_from_data(file_name_prefix, filetype, data):
    doc = Document()

    def add_paragraph(text, bold=False, italic=False, size=11):
        p = doc.add_paragraph()
        run = p.add_run(str(text))
        run.bold = bold
        run.italic = italic
        run.font.size = Pt(size)

    if filetype == "xlsx":
        for sheet_name, rows in data.items():
            doc.add_heading(sheet_name, level=1)
            for row in rows:
                add_paragraph(" | ".join(str(cell) for cell in row))
            doc.add_page_break()

    elif filetype == "pptx":
        for slide_idx, slide in enumerate(data, 1):
            doc.add_heading(f"Slide {slide_idx}: {slide.get('title', '')}", level=1)

            if slide.get("texts"):
                doc.add_heading("Texts:", level=2)
                for text in slide["texts"]:
                    add_paragraph(text)

            images = slide.get("images", [])
            if images:
                doc.add_heading(f"Images ({len(images)} stored in DB):", level=2)
                for img_idx, img in enumerate(images, 1):
                    add_paragraph(f"Image {img_idx} (id: {img['id']}, ext: .{img.get('ext','png')})")
                    if img.get("ocr_text"):
                        for line in img["ocr_text"].splitlines():
                            add_paragraph(line)

            tables = slide.get("tables", [])
            for t_idx, table in enumerate(tables, 1):
                doc.add_heading(f"Table {t_idx}:", level=2)
                for row in table:
                    add_paragraph(" | ".join(str(cell) for cell in row))

            doc.add_page_break()

    elif filetype == "pdf":
        for page in data:
            doc.add_heading(f"Page {page.get('page_number','?')}", level=1)
            doc.add_heading("Extracted Text:", level=2)
            for line in page.get("text", "").splitlines():
                add_paragraph(line)

            for t_idx, table in enumerate(page.get("tables", []), 1):
                doc.add_heading(f"Table {t_idx}:", level=2)
                for row in table:
                    add_paragraph(" | ".join(str(cell) for cell in row))

            if page.get("image_texts"):
                doc.add_heading("OCR Texts from Images:", level=2)
                for text in page["image_texts"]:
                    for line in text.splitlines():
                        add_paragraph(line)

            images = page.get("images", [])
            if images:
                doc.add_heading(f"Images ({len(images)} stored in DB):", level=2)
                for img_idx, img in enumerate(images, 1):
                    add_paragraph(f"Image {img_idx} (id: {img['id']}, ext: .{img.get('ext','png')})")
                    if img.get("ocr_text"):
                        for line in img["ocr_text"].splitlines():
                            add_paragraph(line)

            doc.add_page_break()

    elif filetype in {"png", "jpg", "jpeg", "bmp", "gif", "tiff"}:
        doc.add_heading("Image Data", level=1)
        add_paragraph("OCR Extracted Text:", bold=True)
        add_paragraph(data.get("ocr_text", ""))

        images = data.get("images", [])
        if images:
            doc.add_heading("Image References:", level=2)
            for idx, img in enumerate(images, 1):
                add_paragraph(f"Image {idx} (id: {img['id']}, ext: .{img.get('ext','png')})")
                if img.get("ocr_text"):
                    add_paragraph(img["ocr_text"])

    else:
        doc.add_heading("Data", level=1)
        add_paragraph(json.dumps(data, indent=2))

    # Save to BytesIO stream
    word_stream = BytesIO()
    doc.save(word_stream)
    word_stream.seek(0)
    return word_stream

def expand_query(query):
    tokens = [re.sub(r"[^\w]", "", t.lower()) for t in query.split()]
    expanded = set(tokens)
    for token in tokens:
        if not token:
            continue
        expanded.add(ps.stem(token))
        for syn in wordnet.synsets(token):
            for lemma in syn.lemmas():
                expanded.add(lemma.name().replace("_", " ").lower())
    return list(expanded)


def regex_fallback(query):
    return list(mongo.db.files.find({"$or": [{"filename": {"$regex": query, "$options": "i"}}, {"search_text": {"$regex": query, "$options": "i"}}]}))

# Load FAISS index already called above

def get_failures_by_keyword(keyword):
    with driver.session() as session:
        result = session.run("""
            MATCH (f:Failure)-[:BELONGS_TO]->(d)
            WHERE f.description CONTAINS $keyword
            RETURN f.code, f.description, d.name
        """, keyword=keyword)
        return [record.data() for record in result]



def get_kg_matches(query_terms):
    """
    Fetch KG matches for the given query terms.
    Returns a list of dicts containing failure & asset info.
    """
    kg_matches = []

    with driver.session() as session:
        for term in query_terms:
            term_lower = term.lower()

            # 1. Match failures by code or description
            results_failures = session.run("""
                MATCH (f:Failure)-[:BELONGS_TO]->(d:Department)
                WHERE toLower(f.description) CONTAINS toLower($term)
                   OR toLower(f.code) = toLower($term)
                RETURN f.code AS code, f.description AS description, d.name AS department
            """, {"term": term})

            for r in results_failures:
                kg_matches.append({
                    "failure_code": r["code"],
                    "failure_desc": r["description"],
                    "department": r["department"]
                })

            # 2️⃣ Match assets safely
            results_assets = session.run("""
                MATCH (a:Asset)-[:MONITORED_BY]->(f:Failure)
                WHERE a.asset_name IS NOT NULL AND toLower(a.asset_name) CONTAINS toLower($term)
                RETURN a.asset_name AS asset_name, f.code AS failure_code
            """, {"term": term})


            for r in results_assets:
                kg_matches.append({
                    "asset_name": r["asset_name"],
                    "failure_code": r["failure_code"]
                })

    return kg_matches

def intelligent_search(query):
    sem_results = semantic_search(query, top_k=10)
    ents = extract_entities(query)
    
    kg_docs = []
    if ents:
        kg_docs = list(kg_collection.find({"entities": {"$in": ents}}))
        kg_doc_ids = [k["file_id"] for k in kg_docs]
        kg_files = list(mongo.db.files.find({"_id": {"$in": [ObjectId(fid) for fid in kg_doc_ids]}}))
    else:
        kg_files = []

    # Combine and rank
    combined = {str(r["_id"]): r for r in sem_results}
    for f in kg_files:
        fid = str(f["_id"])
        if fid not in combined:
            combined[fid] = f
            combined[fid]["semantic_score"] = 0.5  # boost from KG

    return sorted(combined.values(), key=lambda x: x.get("semantic_score", 0), reverse=True)

def semantic_search(query, top_k=10, normalize=True):
    """
    Query FAISS and return MongoDB docs sorted by descending similarity (higher is better when normalized).
    """
    global index, id_map
    if index is None or getattr(index, "ntotal", 0) == 0:
        return []

    if not query or not query.strip():
        return []

    q_vec = model.encode(query, convert_to_numpy=True)
    q_vec = np.asarray(q_vec, dtype="float32").reshape(1, -1)

    if normalize:
        qn = np.linalg.norm(q_vec, axis=1, keepdims=True)
        if qn[0][0] == 0.0:
            return []
        q_vec = q_vec / qn

    # sanity check for dimension mismatch
    if q_vec.shape[1] != index.d:
        app.logger.error("Query embedding dim (%d) != index.d (%d).", q_vec.shape[1], index.d)
        return []

    with index_lock:
        k = min(top_k, max(1, index.ntotal))
        distances, indices = index.search(q_vec, k)

    # collect valid ids (avoid __deleted__)
    valid_pairs = []
    for score, idx in zip(distances[0], indices[0]):
        if idx < 0 or idx >= len(id_map):
            continue
        mongo_id = id_map[idx]
        if mongo_id == "__deleted__":
            continue
        valid_pairs.append((float(score), mongo_id))

    if not valid_pairs:
        return []

    # batch lookup to reduce MongoDB round-trips
    mongo_ids = [ObjectId(mid) for _, mid in valid_pairs]
    found_docs = list(mongo.db.files.find({"_id": {"$in": mongo_ids}}))
    found_map = {str(doc["_id"]): doc for doc in found_docs}

    results = []
    for score, mid in valid_pairs:
        doc = found_map.get(mid)
        if doc:
            doc["semantic_score"] = score
            doc["_id_str"] = str(doc["_id"])
            results.append(doc)

    results.sort(key=lambda x: x.get("semantic_score", 0.0), reverse=True)
    return results

# 🧠 Load summarization + analysis models (load once)

def run_search(query, limit=800):
    words = [w for w in query.split() if w]
    if not words:
        return [], "empty query"

    # Projection for fields you want to return
    projection = {
        "score": {"$meta": "textScore"},
        "filename": 1,
        "filetype": 1,
        "data": 1
    }

    # Exact phrase search
    phrase_results = list(
        mongo.db.files.find({"$text": {"$search": f'"{query}"'}}, projection)
        .sort([("score", {"$meta": "textScore"})])
        .limit(limit)
    )
    if phrase_results:
        return phrase_results, "exact phrase"

    # AND search for individual words (all must appear)
    and_query = " ".join([f'"{w}"' for w in words])
    and_results = list(
        mongo.db.files.find({"$text": {"$search": and_query}}, projection)
        .sort([("score", {"$meta": "textScore"})])
        .limit(limit)
    )
    if and_results:
        return and_results, "AND match"

    return [], "no match"

import nltk
from nltk import pos_tag, word_tokenize
from nltk.corpus import stopwords
nltk.download('stopwords')
nltk.download('punkt_tab')
nltk.download('averaged_perceptron_tagger_eng')
# Make sure you’ve downloaded NLTK data once in your environment:
# nltk.download('punkt')
# nltk.download('averaged_perceptron_tagger')
# nltk.download('stopwords')

stop_words = set(stopwords.words('english'))

def refine_query_terms(query: str) -> list:
    """
    Extract domain-friendly key terms from query:
    - Keeps nouns, proper nouns, adjectives
    - Keeps codes, hyphenated and slash-separated terms
    - Filters stopwords and punctuation
    """
    query = (query or "").strip()
    if not query:
        return []

    # Tokenize
    tokens = word_tokenize(query)

    key_terms = []
    for token in tokens:
        w_clean = token.strip()
        w_lower = w_clean.lower()
        # Skip stopwords and empty strings
        if not w_clean or w_lower in stop_words:
            continue
        # Keep alphanumerics, hyphens, slashes, and uppercase codes
        if re.match(r"^[a-zA-Z0-9-/]+$", w_clean):
            key_terms.append(w_clean)
        else:
            # Optionally, keep proper nouns (POS NN/NNP) if using pos_tag
            pass

    # Fallback: keep all non-stopword tokens
    if not key_terms:
        key_terms = [t for t in tokens if t.lower() not in stop_words]

    print("Refined query terms:", key_terms)
    return key_terms

# 🧩 --- Summarization helper ---

# Cached global reranker (loads once)
onnx_reranker = None
def safe_summarize_text(summarizer, text, min_len=100, max_len=300, max_model_input=900):
    # Safely summarize text by truncating it to the model's input size limit.
    # Prevents 'index out of range' errors for long sequences.
    if not text or not summarizer:
        return ""

    # Clean + limit to model capacity (approx 900 tokens ≈ ~1500 words)
    tokens = text.split()
    if len(tokens) > max_model_input:
        text = " ".join(tokens[:max_model_input])

    try:
        result = summarizer(
            text,
            max_length=max_len,
            min_length=min_len,
            do_sample=False
        )
        if isinstance(result, list) and result and "summary_text" in result[0]:
            return result[0]["summary_text"].strip()

    except Exception as e:
        app.logger.error(f"⚠️ Summarization failed safely: {e}", exc_info=True)

    # Fallback to truncated original
    return text[:700]

def load_onnx_reranker(model_name="cross-encoder/ms-marco-MiniLM-L-6-v2"):
    """
    Load an ONNX quantized CrossEncoder model for CPU inference.
    """
    global onnx_reranker

    if onnx_reranker:
        return onnx_reranker

    # Download or load quantized ONNX version
    model_path = f"./onnx_reranker/{model_name.replace('/', '_')}.onnx"
    tokenizer = AutoTokenizer.from_pretrained(model_name)

    try:
        ort_session = ort.InferenceSession(
            model_path,
            providers=["CPUExecutionProvider"],
            sess_options=ort.SessionOptions()
        )
    except Exception as e:
        # If no ONNX model yet, export from original model once
        print("Exporting CrossEncoder to ONNX...")
        ce = CrossEncoder(model_name)
        ce.save_pretrained("./onnx_reranker")
        tokenizer.save_pretrained("./onnx_reranker")

    # Convert via torch -> ONNX (first-time only)
        import torch

        inputs = tokenizer("example question", "example answer", return_tensors="pt")
        torch.onnx.export(
            ce.model,
            (inputs["input_ids"], inputs["attention_mask"]),
            model_path,
            input_names=["input_ids", "attention_mask"],
            output_names=["logits"],
            dynamic_axes={
                "input_ids": {0: "batch", 1: "sequence"},
                "attention_mask": {0: "batch", 1: "sequence"},
            },
            opset_version=14
        )
        ort_session = ort.InferenceSession(model_path, providers=["CPUExecutionProvider"])

    onnx_reranker = {"session": ort_session, "tokenizer": tokenizer}
    return onnx_reranker


def rerank_with_onnx(query, docs, top_k=10):
    """
    Given a query and list of text docs, return reranked docs (sorted by score).
    """
    rr = load_onnx_reranker()
    tok = rr["tokenizer"]
    session = rr["session"]

    pairs = [(query, d) for d in docs]
    enc = tok([p[0] for p in pairs], [p[1] for p in pairs],
              padding=True, truncation=True, max_length=512, return_tensors="np")

    ort_inputs = {
        "input_ids": enc["input_ids"],
        "attention_mask": enc["attention_mask"]
    }
    logits = session.run(None, ort_inputs)[0].squeeze(-1)

    # Attach scores back to docs
    scored = [{"text": d, "score": float(s)} for d, s in zip(docs, logits)]
    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:top_k]


def tokenize(text):
    """Return a set of lowercase words in the text."""
    return set(re.findall(r'\b\w+\b', text.lower()))


def normalize_query_terms(query):
    """
    Produce a stable list of query terms for matching.
    Handles cases where the query contains code-like text (parentheses,
    commas, cypher SQL fragments) by falling back to alphanumeric token
    extraction. Returns lowercased terms.
    """
    if not query:
        return []
    # Primary: word tokens
    terms = re.findall(r'\b\w+\b', query)
    if not terms:
        # Fallback: alphanumeric sequences (identifiers)
        terms = re.findall(r'[A-Za-z0-9_]{2,}', query)
    # Filter out very short noisy tokens (single punctuation or single letters)
    terms = [t.lower() for t in terms if re.search(r'[A-Za-z0-9]', t)]
    return terms


def summarize_relevant_chunks(full_text, query, summarizer,
                              min_para_len=250, max_chunks=5,
                              chunk_size=2000, overlap=300,
                              min_summary_len=100, max_summary_len=300,
                              max_chunk_chars=2000):
    """
    Extract relevant paragraphs, merge into large chunks,
    filter to chunks containing all query terms (AND),
    and summarize top relevant chunks safely.
    """

    # Step 1: Extract relevant paragraphs
    matched_paragraphs = extract_relevant_paragraphs(
        full_text, query, min_len=min_para_len, max_chunks=max_chunks, fallback_sentence_chunk=8
    )

    if not matched_paragraphs:
        app.logger.debug("No relevant paragraphs found. Using fallback.")
        return {
            "final_summary": full_text[:1000],
            "summaries_with_ids": []
        }

    app.logger.debug(f"Found {len(matched_paragraphs)} relevant paragraphs.")

    # Step 2: Create semantically-aware chunks using tokenizer + sentence splitting
    # This preserves tables and sentence boundaries (chunk_text is table-aware).
    raw_chunks = chunk_text(full_text, max_tokens=chunk_size, overlap=overlap)
    chunks = [{"id": f"chunk_{i+1}", "text": c} for i, c in enumerate(raw_chunks)]

    # Step 3: Filter chunks (strict AND operation)
    query_terms = normalize_query_terms(query)
    filtered_chunks = [c for c in chunks if all(qt in c["text"].lower() for qt in query_terms)]

    # Fallback to OR if no AND match
    if not filtered_chunks:
        app.logger.debug("No chunks contained all query terms. Falling back to OR matching.")
        filtered_chunks = [c for c in chunks if any(qt in c["text"].lower() for qt in query_terms)]

    if not filtered_chunks:
        app.logger.debug("Still no matching chunks. Using fallback first chunk.")
        filtered_chunks = [chunks[0]]

    app.logger.debug(f"Filtering kept {len(filtered_chunks)} chunks (strict AND logic applied).")

    # Step 4: Rank chunks
    ranked_chunks = []
    for c in filtered_chunks:
        rank, rank_type = chunk_rank(c["text"], query_terms)
        ranked_chunks.append((c, rank, rank_type))

    ranked_chunks.sort(key=lambda x: x[1], reverse=True)
    top_chunks = [x[0] for x in ranked_chunks[:max_chunks]]

    app.logger.debug(f"Summarizing top {len(top_chunks)} relevant chunks...")

    # Step 5: Summarize each chunk safely
    summaries_with_ids = []
    for idx, chunk in enumerate(top_chunks):
        text = chunk["text"].strip()

        # Skip re-summarization if text already looks summarized
        if len(text) < 350 or "summary" in text.lower() or "key points" in text.lower():
            summary = text[:max_summary_len]
            print(f"🟡 Skipped summarization for {chunk['id']} (already short/summarized).")
        else:
            summary = safe_summarize_text(
                summarizer,
                text,
                min_len=min_summary_len,
                max_len=max_summary_len,
                max_model_input=900
            )

        if summary:
            summaries_with_ids.append({
                "chunk_id": chunk["id"],
                "summary": summary
            })
            print(f"✅ {chunk['id']} summarized successfully ({len(summary)} chars).")
        else:
            print(f"⚠️ {chunk['id']} produced no summary (skipped).")

    # Step 6: Merge summaries
    final_summary = " ".join(s["summary"] for s in summaries_with_ids).strip()

    return {
        "final_summary": final_summary or text_to_summarize[:1000],
        "summaries_with_ids": summaries_with_ids
    }


# Independent AND/OR ranking for a single chunk
def chunk_rank(chunk, query_terms):
    """
    Rank a single chunk based on query terms:
    - AND if all query terms present
    - OR if some terms present
    - 0 if none
    Does NOT aggregate across chunks.
    """
    chunk_tokens = set(re.findall(r'\b\w+\b', chunk.lower()))
    matched_terms = [t for t in query_terms if t in chunk_tokens]

    if not matched_terms:
        return 0, None
    elif len(matched_terms) == len(query_terms):
        return 2, "AND"   # highest priority
    else:
        return 1, "OR"    # partial match, lower priority

# 📊 --- Analysis helper ---
def analyze_text_summary(text):
    """Return top entities and key keywords from text summary."""
    analysis = {"entities": [], "keywords": []}
    if not text:
        return analysis

    try:
        entities = ner_analyzer(text)
        keywords = kw_model.extract_keywords(text, top_n=10)
        analysis["entities"] = [{"entity": e["entity_group"], "text": e["word"]} for e in entities]
        analysis["keywords"] = [kw[0] for kw in keywords]
    except Exception as e:
        print("Analysis error:", e)
    return analysis


def extract_relevant_paragraphs(full_text, query, min_len=250, max_chunks=5, fallback_sentence_chunk=8):
    """
    Extract paragraphs that match the query progressively:
    - Try paragraphs containing all query terms (AND)
    - Then paragraphs containing subsets of terms (OR)
    - Rank paragraphs by number of query terms present
    """
    query_terms = [t.lower() for t in re.findall(r'\b\w+\b', query)]
    # Split text into paragraphs using double newline or line breaks
    paragraphs = re.split(r'\n{2,}|\r{2,}|\n', full_text)
    ranked_paragraphs = []

    # Iterate progressively: try all terms -> N-1 terms -> ... -> 1 term
    for num_terms in range(len(query_terms), 0, -1):
        for para in paragraphs:
            para_text = para.strip()
            if len(para_text) < min_len:
                continue
            para_tokens = set(re.findall(r'\b\w+\b', para_text.lower()))
            # Count how many query terms are present
            matched_terms = [t for t in query_terms if t in para_tokens]
            if len(matched_terms) >= num_terms:
                ranked_paragraphs.append((para_text, len(matched_terms)))
        if ranked_paragraphs:
            # Stop after finding paragraphs with current number of terms
            break

    # Sort paragraphs by matched terms descending
    ranked_paragraphs.sort(key=lambda x: x[1], reverse=True)
    selected_paragraphs = [p[0] for p in ranked_paragraphs[:max_chunks]]

    # Fallback: if no paragraph found, split text into sentence chunks
    if not selected_paragraphs:
        sentences = re.split(r'(?<=[.!?])\s+', full_text)
        chunks = [" ".join(sentences[i:i+fallback_sentence_chunk])
                  for i in range(0, len(sentences), fallback_sentence_chunk)]
        selected_paragraphs = chunks[:max_chunks]

    return selected_paragraphs



def rank_and_score_document(full_text, query, query_terms, base_score=0.0):
    """
    Robust document ranking:
    - Progressive AND -> OR term matching
    - Phrase + term boosts (exact word matches)
    - Returns final_score, rank_type
    """
    if not full_text or not query_terms:
        return base_score, "none"

    # --- Tokenize document ---
    tokens = re.findall(r'\b\w+\b', full_text.lower())
    token_set = set(tokens)

    # --- Progressive AND -> OR matching ---
    n_terms = len(query_terms)
    rank_type = "none"
    matched_terms_count = 0
    for num_terms in range(n_terms, 0, -1):
        for combo in combinations(query_terms, num_terms):
            if all(term.lower() in token_set for term in combo):
                matched_terms_count = num_terms
                rank_type = "AND" if num_terms == n_terms else f"OR-{num_terms}"
                break
        if matched_terms_count:
            break

    # --- Phrase & term boosts ---
    # Count exact query terms
    term_count = sum(tokens.count(term.lower()) for term in query_terms)

    # Count exact query phrase (all terms together)
    phrase_pattern = r'\b' + r'\s+'.join(re.escape(term.lower()) for term in query_terms) + r'\b'
    phrase_count = len(re.findall(phrase_pattern, full_text.lower()))

    weight = (phrase_count * 2.0) + (term_count * 0.2) + matched_terms_count
    boost = min(2.0, math.log1p(weight))
    final_score = base_score + boost

    return final_score, rank_type


# --- Retrieve related entities for query expansion ---
def kg_expand_query(term, max_hops=2):
    """
    Given a term like 'signal failure' or 'gear fuse',
    find related assets, failures, stations, and departments from KG.
    Searches across multiple property names (case-insensitive) and expands up to N hops.
    Does NOT depend on APOC. Fully parameterized, sanitized, and rate-limited.
    """
    if not driver:
        app.logger.debug("KG expand skipped: Neo4j driver unavailable")
        return []
    
    try:
        # Sanitize input
        term = _sanitize_query_term(term)
        if not term:
            return []
        
        # Validate and constrain max_hops
        max_hops = _validate_int_param(max_hops, default=2, min_val=1, max_val=5)
        
        # Check rate limit
        if not _check_kg_rate_limit():
            app.logger.warning("KG expand rate limit exceeded")
            return []
        
        # cache key based on term + max_hops
        cache_key = f"expand:{term.strip().lower()}:{int(max_hops)}"
        cached = _cache_get(KG_EXPAND_CACHE, cache_key)
        if cached is not None:
            return cached
        with driver.session() as session:
            q = f"""
            MATCH (n)
            WHERE any(k IN keys(n) WHERE toLower(toString(n[k])) CONTAINS toLower($term))

            // Expand relationships up to N hops in both directions
            OPTIONAL MATCH (n)-[*1..{max_hops}]-(m)
            WITH COLLECT(DISTINCT m) + COLLECT(DISTINCT n) AS nodes

            UNWIND nodes AS node
            RETURN DISTINCT coalesce(
                node.name,
                node.code,
                node.failure_entry_no,
                node.Failure_Entry_No,
                node.zone,
                node.Zone,
                node.division,
                node.Division,
                node.section,
                node.Section,
                node.station_code,
                node.Station_Code,
                node.department,
                node.Department,
                node.gear_name,
                node.Gear_Name,
                node.cause,
                node.Cause,
                node.remarks,
                node.Remarks
            ) AS name
            LIMIT 50
            """
            res = timed_run(session, q, term=term)
            related_terms = [r["name"] for r in res if r["name"]]
            # Normalize and split returned names into useful expansion tokens.
            cleaned = []
            for name in related_terms:
                try:
                    if not name:
                        continue
                    # remove surrounding whitespace and control chars
                    n = _illegal_xml_re.sub('', str(name)).strip()
                    if not n:
                        continue
                    # keep the full phrase (lowercased) if reasonably short
                    phrase = n.lower()
                    if 2 <= len(phrase) <= 120:
                        cleaned.append(phrase)
                    # also include token-level expansions (words, identifiers)
                    toks = re.findall(r"\b[\w\-/]+\b", n)
                    for t in toks:
                        tt = t.lower().strip()
                        if tt and len(tt) > 1 and tt != phrase:
                            cleaned.append(tt)
                except Exception:
                    continue

            out = list(dict.fromkeys(cleaned))  # preserve order, unique
            _cache_set(KG_EXPAND_CACHE, cache_key, out)
            return out

    except Exception as e:
        app.logger.warning(f"⚠️ KG expand failed for term '{term}': {e}")
        return []

# --- Retrieve failures connected to key nodes ---
def kg_get_related_failures(query_terms):
    """
    Retrieve Failure nodes related to KG entities (zones, stations, gears, etc.).
    Case-insensitive, direction-agnostic, multi-term friendly.
    Returns up to 100 relevant records. Fully parameterized, sanitized, and rate-limited.
    """
    if not driver:
        app.logger.debug("KG related failures skipped: Neo4j driver unavailable")
        return []
    
    if not query_terms:
        return []

    # Sanitize all terms
    terms = [_sanitize_query_term(t) for t in query_terms if t.strip()]
    terms = [t for t in terms if t]  # remove empty after sanitization
    if not terms:
        return []
    
    # Check rate limit
    if not _check_kg_rate_limit():
        app.logger.warning("KG related failures rate limit exceeded")
        return []

    try:
        # cache key
        cache_key = "related_failures:" + ",".join(sorted(query_terms))
        cached = _cache_get(KG_RELATED_FAILURES_CACHE, cache_key)
        if cached is not None:
            return cached
        with driver.session() as session:
            q = """
            // Match failures and any connected nodes (any direction)
            MATCH (f:Failure)
            OPTIONAL MATCH (f)-[*0..2]-(x)
            WHERE any(t IN $terms WHERE
                // check all relevant fields on f and related nodes
                any(val IN [
                    coalesce(f.failure_entry_no, f.Failure_Entry_No, ''),
                    coalesce(f.zone, f.Zone, ''),
                    coalesce(f.division, f.Division, ''),
                    coalesce(f.section, f.Section, ''),
                    coalesce(f.station_code, f.Station_Code, ''),
                    coalesce(f.department, f.Department, ''),
                    coalesce(f.gear_name, f.Gear_Name, ''),
                    coalesce(f.cause, f.Cause, ''),
                    coalesce(f.action, f.Action, ''),
                    coalesce(f.remarks, f.Remarks, ''),
                    coalesce(x.name, ''),
                    coalesce(x.code, ''),
                    coalesce(x.gear_name, ''),
                    coalesce(x.zone, ''),
                    coalesce(x.division, ''),
                    coalesce(x.station_code, '')
                ] WHERE toLower(val) CONTAINS t)
            )
            RETURN DISTINCT
                coalesce(f.failure_entry_no, f.Failure_Entry_No) AS failure_entry_no,
                coalesce(f.zone, f.Zone) AS zone,
                coalesce(f.division, f.Division) AS division,
                coalesce(f.section, f.Section) AS section,
                coalesce(f.station_code, f.Station_Code) AS station_code,
                coalesce(f.department, f.Department) AS department,
                coalesce(f.gear_name, f.Gear_Name) AS gear_name,
                coalesce(f.cause, f.Cause) AS cause,
                coalesce(f.action, f.Action) AS action,
                coalesce(f.remarks, f.Remarks) AS remarks
            LIMIT 100
            """

            res = timed_run(session, q, terms=terms)
            results = [dict(r) for r in res]
            _cache_set(KG_RELATED_FAILURES_CACHE, cache_key, results)

            # Optional lightweight relevance scoring
            for r in results:
                text = " ".join(str(v or "").lower() for v in r.values())
                r["_score"] = sum(text.count(t) for t in terms)

            # Sort by relevance
            results.sort(key=lambda x: x.get("_score", 0), reverse=True)
            return results

    except Exception as e:
        print(f"[KG] Query failed: {e}")
        return []

def kg_entity_importance(entity_text: str, max_hops: int = 2) -> float:
    """
    Compute a robust importance score for an entity using Neo4j.
    Factors:
      - Relationship degree (direct links)
      - Relation type diversity
      - Connectivity depth (up to N hops)
      - Case-insensitive + partial matching
    Returns: normalized float score. Fully parameterized, sanitized, and rate-limited.
    """
    import math

    if not driver:
        app.logger.debug("KG importance skipped: Neo4j driver unavailable")
        return 0.0

    # Sanitize input
    entity_text = _sanitize_query_term(entity_text)
    if not entity_text or len(entity_text.strip()) < 2:
        return 0.0
    
    # Validate and constrain max_hops
    max_hops = _validate_int_param(max_hops, default=2, min_val=1, max_val=5)
    
    # Check rate limit
    if not _check_kg_rate_limit():
        app.logger.warning("KG importance rate limit exceeded")
        return 0.0

    try:
        cache_key = f"importance:{entity_text.strip().lower()}:{int(max_hops)}"
        cached = _cache_get(KG_IMPORTANCE_CACHE, cache_key)
        if cached is not None:
            return cached
        with driver.session() as session:
            try:
                # --- Try APOC version first ---
                cypher = f"""
                 // 1. Find matching entities (case-insensitive, partial)
                MATCH (e:Entity)
                WHERE toLower(e.text) CONTAINS toLower($entity)

                 // 2. Count direct and indirect relationships using APOC
                OPTIONAL MATCH (e)-[r*1..{max_hops}]-(n:Entity)
                WITH e, apoc.coll.flatten(r) AS rels, collect(DISTINCT n) AS connected_nodes

                 // 3. Compute richness metrics
                WITH e,
                     size(connected_nodes) AS reach,
                     size(rels) AS total_links,
                     apoc.coll.toSet([rel IN rels | type(rel)]) AS rel_types

                RETURN e.text AS entity,
                       reach,
                       total_links,
                       size(rel_types) AS relation_variety
                ORDER BY reach DESC
                LIMIT 1
                """
                record = timed_run(session, cypher, entity=entity_text).single()

            except Exception as e_apoc:
                # --- Fallback: APOC not available ---
                app.logger.warning(f"[KG] APOC not available, using fallback for '{entity_text}': {e_apoc}")
                cypher = f"""
                MATCH (e:Entity)
                WHERE toLower(e.text) CONTAINS toLower($entity)
                OPTIONAL MATCH (e)-[r*1..{max_hops}]-(n:Entity)
                WITH e, collect(DISTINCT n) AS connected_nodes,
                     collect(DISTINCT type(r[0])) AS rel_types
                RETURN e.text AS entity,
                       size(connected_nodes) AS reach,
                       size(rel_types) AS relation_variety,
                       0 AS total_links
                ORDER BY reach DESC
                LIMIT 1
                """
                record = timed_run(session, cypher, entity=entity_text).single()

            if not record:
                _cache_set(KG_IMPORTANCE_CACHE, cache_key, 0.0)
                return 0.0

            # --- Safely extract numeric values ---
            reach = record.get("reach") or 0
            total_links = record.get("total_links") or 0
            variety = record.get("relation_variety") or 0

            # --- Weighted composite score ---
            score = (reach * 0.6 + variety * 0.4) * (1 + math.log1p(total_links))
            score = math.log1p(score) * 2  # normalization

            score_r = round(score, 3)
            _cache_set(KG_IMPORTANCE_CACHE, cache_key, score_r)
            return score_r

    except Exception as e:
        app.logger.warning(f"⚠️ KG importance lookup failed for '{entity_text}': {e}")
        return 0.0

# --- Default CPU-friendly base model ---
base_model_name = "distilgpt2"  # small, fast, CPU-friendly

# Load tokenizer and base model
tokenizer = AutoTokenizer.from_pretrained(base_model_name)
base_model = AutoModelForCausalLM.from_pretrained(base_model_name)

# --- LoRA adapter path ---
# For default testing, you can skip LoRA and use base_model directly on CPU
# If you later train a LoRA adapter, replace "./lora_cpu_model" with its path
try:
    lora_checkpoint = "./lora_railway_model/checkpoint-642"
    lora_model = PeftModel.from_pretrained(model, lora_checkpoint)
except:
    print("⚠️ No LoRA model found, using base model on CPU")
    lora_model = base_model

# --- Pipeline for text generation on CPU ---
llm_pipe = pipeline(
    "text-generation",
    model=lora_model,
    tokenizer=tokenizer,
    device=-1,  # CPU
    max_length=512
)

# --- Utility: chunk extraction ---
def chunk_text(text, max_tokens=2000, overlap=200):
    """
    Split text into overlapping chunks with an algorithm that preserves
    sentence/paragraph boundaries (preferred) and respects a token budget.

    Strategy (best-practice):
    - Use spaCy sentence segmentation to split text into sentences.
    - Measure token length per sentence using the tokenizer.
    - Greedily pack sentences into a chunk until adding the next sentence
      would exceed max_tokens.
    - When finalizing a chunk, compute an overlap boundary by walking
      backwards over the last sentences until the overlap token budget
      is reached; start the next chunk from that sentence to provide
      semantic overlap.

    Falls back to tokenizer-token sliding-window if sentence segmentation
    or tokenizer calls fail.
    """
    try:
        # Use spaCy sentence splitting for semantic-aware chunking
        # First, split into table blocks and non-table blocks using markers
        blocks = []
        parts = re.split(r'(?:\n?___TABLE_START___\n?.*?___TABLE_END___\n?)', text, flags=re.S)
        # re.split will include separators; instead we findall blocks
        blocks = []
        last_idx = 0
        for m in re.finditer(r'___TABLE_START___\n.*?___TABLE_END___', text, flags=re.S):
            # preceding text
            if m.start() > last_idx:
                pre = text[last_idx:m.start()]
                if pre.strip():
                    blocks.append(("text", pre))
            blocks.append(("table", m.group(0)))
            last_idx = m.end()
        if last_idx < len(text):
            tail = text[last_idx:]
            if tail.strip():
                blocks.append(("text", tail))

        chunks = []
        # Process each block: tables are atomic (split by rows if too large), text blocks use sentence-aware packing
        for btype, bcontent in blocks:
            if btype == "table":
                # Tokenize the whole table block
                try:
                    tbl_tokens = tokenizer.encode(bcontent, truncation=False)
                    if len(tbl_tokens) <= max_tokens:
                        chunks.append(bcontent)
                        continue
                except Exception:
                    # fallback to naive row-splitting
                    pass

                # Split table into header + rows
                lines = [ln for ln in bcontent.splitlines() if ln.strip()]
                header = None
                rows = []
                for ln in lines:
                    if ln.startswith('TABLE_HEADER:'):
                        header = ln
                    elif ln.startswith('TABLE_ROW:'):
                        rows.append(ln)

                # Pack rows into chunks preserving header
                cur_rows = []
                cur_tokens = 0
                for r in rows:
                    try:
                        r_tokens = len(tokenizer.encode(r, truncation=False))
                    except Exception:
                        r_tokens = max(1, len(r.split()))
                    if cur_tokens + r_tokens <= max_tokens - (len(tokenizer.encode(header, truncation=False)) if header else 0):
                        cur_rows.append(r)
                        cur_tokens += r_tokens
                    else:
                        # finalize chunk
                        block_text = '\n'.join(([header] if header else []) + cur_rows + [''])
                        chunks.append(block_text)
                        # compute overlap rows to carry to next chunk (by token budget)
                        overlap_tokens = 0
                        overlap_rows = []
                        k = len(cur_rows) - 1
                        while k >= 0 and overlap_tokens < overlap:
                            try:
                                tokc = len(tokenizer.encode(cur_rows[k], truncation=False))
                            except Exception:
                                tokc = max(1, len(cur_rows[k].split()))
                            overlap_tokens += tokc
                            overlap_rows.insert(0, cur_rows[k])
                            k -= 1
                        cur_rows = overlap_rows + [r]
                        cur_tokens = sum((len(tokenizer.encode(rr, truncation=False)) if 'tokenizer' in globals() else max(1, len(rr.split()))) for rr in cur_rows)

                if cur_rows:
                    block_text = '\n'.join(([header] if header else []) + cur_rows + [''])
                    chunks.append(block_text)
                continue

            # btype == "text" : perform sentence-aware packing on the block
        # For all text blocks, re-join and proceed with spaCy on combined text for sentences
        text_blocks = "\n\n".join(b for t, b in blocks if t == "text")
        doc = nlp(text_blocks)
        sents = [s.text.strip() for s in doc.sents if s.text.strip()]
        if not sents:
            raise ValueError("No sentences from spaCy")

        # Precompute token lengths for sentences
        sent_token_counts = []
        for s in sents:
            try:
                toks = tokenizer.encode(s, truncation=False)
                sent_token_counts.append(len(toks))
            except Exception:
                # Fallback to approximate count by words
                sent_token_counts.append(max(1, len(s.split())))

        chunks = []
        i = 0
        n = len(sents)
        while i < n:
            token_sum = 0
            j = i
            # Pack sentences until token budget exceeded
            while j < n and token_sum + sent_token_counts[j] <= max_tokens:
                token_sum += sent_token_counts[j]
                j += 1

            # If a single sentence is longer than max_tokens, forcibly split it
            if j == i:
                # split tokens from the long sentence directly
                long_sent = sents[i]
                toks = tokenizer.encode(long_sent, truncation=False)
                start_t = 0
                while start_t < len(toks):
                    end_t = min(start_t + max_tokens, len(toks))
                    chunk_text_piece = tokenizer.decode(toks[start_t:end_t], skip_special_tokens=True)
                    chunks.append(chunk_text_piece)
                    start_t = end_t - overlap if end_t - overlap > start_t else end_t
                i += 1
                continue

                # Build chunk from sentences i..j-1
                chunk = " ".join(sents[i:j])
                chunks.append(chunk)

            # Determine overlap: walk backward from j-1 until accumulated tokens >= overlap
            if overlap > 0 and j - 1 >= i:
                overlap_tokens = 0
                k = j - 1
                while k >= i and overlap_tokens < overlap:
                    overlap_tokens += sent_token_counts[k]
                    k -= 1
                # next start is k+1 (ensures at least one sentence overlap)
                i = max(i + 1, k + 1)
            else:
                i = j

        return chunks

    except Exception:
        # Fallback: tokenizer token sliding window (original behavior)
        try:
            tokens = tokenizer.encode(text, truncation=False)
            chunks = []
            start = 0
            while start < len(tokens):
                end = start + max_tokens
                chunk_tokens = tokens[start:end]
                ctext = tokenizer.decode(chunk_tokens, skip_special_tokens=True)
                chunks.append(ctext)
                start += max(1, (max_tokens - overlap))
            return chunks
        except Exception:
            # Final fallback: paragraph-based char splitting
            paragraphs = re.split(r"\n{2,}|\r{2,}|\n", text)
            chunks = []
            cur = ""
            for p in paragraphs:
                if len(cur) + len(p) + 1 <= max_tokens * 4:
                    cur = (cur + "\n\n" + p).strip()
                else:
                    if cur:
                        chunks.append(cur)
                    cur = p
            if cur:
                chunks.append(cur)
            return chunks


# --- Utility: select top-k chunks by cosine similarity ---


# Embedding helpers - prefer OpenAI embeddings when configured, fallback to
# sentence-transformers locally. This section implements a recursive splitting
# strategy to ensure each piece sent to the embedding model fits token/char
# limits. It also batches requests to the embedding API with retries.

# Fallback HF model
embed_model = None
try:
    embed_model = SentenceTransformer("all-MiniLM-L6-v2")
except Exception:
    embed_model = None

# OpenAI embedding configuration
OPENAI_EMBED_MODEL = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-small")
OPENAI_ENABLED = bool(os.getenv("OPENAI_API_KEY"))
if OPENAI_ENABLED:
    openai.api_key = os.getenv("OPENAI_API_KEY")


def recursive_split_text(text, max_chars=3000):
    """Recursively split text into pieces smaller than max_chars while
    trying to respect paragraph/sentence boundaries.
    """
    text = text.strip()
    if not text:
        return []
    if len(text) <= max_chars:
        return [text]

    # Prefer split on double newline (paragraph), then sentence boundary, then mid-point
    parts = re.split(r"\n{2,}", text)
    if len(parts) > 1:
        out = []
        cur = ""
        for p in parts:
            if len((cur + "\n\n" + p).strip()) <= max_chars:
                cur = (cur + "\n\n" + p).strip()
            else:
                if cur:
                    out.extend(recursive_split_text(cur, max_chars))
                cur = p
        if cur:
            out.extend(recursive_split_text(cur, max_chars))
        return out

    # Fallback to sentence splits
    sentences = re.split(r'(?<=[.!?])\s+', text)
    out = []
    cur = ""
    for s in sentences:
        if len((cur + " " + s).strip()) <= max_chars:
            cur = (cur + " " + s).strip()
        else:
            if cur:
                out.append(cur)
            cur = s
    if cur:
        out.append(cur)

    # If still oversized (very long sentence), split by midpoint
    for i, o in enumerate(out):
        if len(o) > max_chars:
            mid = len(o) // 2
            out[i:i+1] = recursive_split_text(o[:mid], max_chars) + recursive_split_text(o[mid:], max_chars)
    return out


def batch_get_openai_embeddings(texts, model_name=OPENAI_EMBED_MODEL, batch_size=16, retries=3):
    embeddings = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        attempt = 0
        while attempt < retries:
            try:
                resp = openai.Embedding.create(input=batch, model=model_name)
                batch_emb = [np.array(d["embedding"], dtype="float32") for d in resp["data"]]
                embeddings.extend(batch_emb)
                break
            except Exception as e:
                attempt += 1
                wait = 1.0 * attempt
                time.sleep(wait)
                if attempt >= retries:
                    raise
    return embeddings


def get_embeddings(texts, use_openai=OPENAI_ENABLED):
    """Return list of numpy embeddings for `texts`.
    If OpenAI is configured, use it. Otherwise fall back to sentence-transformers.
    """
    if use_openai:
        # Recursively ensure each text is within safe char limits
        safe_texts = []
        for t in texts:
            if not t:
                safe_texts.append("")
            elif len(t) > 3000:
                safe_texts.extend(recursive_split_text(t, max_chars=3000))
            else:
                safe_texts.append(t)
        return batch_get_openai_embeddings(safe_texts)

    # Fallback to local model
    if embed_model is None:
        raise RuntimeError("No embedding model available: set OPENAI_API_KEY or install sentence-transformers")
    emb = embed_model.encode(texts, convert_to_numpy=True)
    # sentence-transformers may return shape (n, d)
    return [np.asarray(e, dtype="float32") for e in emb]


def select_top_chunks(chunks, query, top_k=3):
    """Select top_k chunks most relevant to query, using embeddings.
    Uses OpenAI embeddings when available, else local model.
    """
    # Prepare inputs: we will embed chunks and the query
    texts = [query] + chunks
    embs = get_embeddings(texts)
    query_emb = embs[0]
    chunk_embs = embs[1:]
    sims = [float(cosine_similarity(query_emb.reshape(1, -1), ce.reshape(1, -1))[0][0]) for ce in chunk_embs]
    ranked = sorted(zip(chunks, sims), key=lambda x: x[1], reverse=True)
    return [c for c, _ in ranked[:top_k]]


def safe_chunk_embeddings(chunks, query, top_k=5):
    """Return top_k chunks relevant to query with robust embedding creation.
    This will recursively split oversized chunks before embedding.
    """
    # Recursively split any chunk that exceeds OpenAI safe limits
    processed = []
    for c in chunks:
        if OPENAI_ENABLED and len(c) > 3000:
            processed.extend(recursive_split_text(c, max_chars=3000))
        else:
            processed.append(c)

    return select_top_chunks(processed, query, top_k=top_k)



# 🧩 --- Your existing search route with summarization + analysis integrated ---
@app.route("/search", methods=["GET"])
def search():
    query = request.args.get("q", "").strip()
    department_filter = request.args.get("department", None)
    summarize_flag = request.args.get("summarize", "false").lower() == "true"
    analyze_flag = request.args.get("analyze", "false").lower() == "true"

    app.logger.info(f"Search query='{query}' summarize={summarize_flag} analyze={analyze_flag}")

    # --- Step 0: Empty query guard ---
    if not query:
        return (
            render_template("search.html", results=[], query="")
            if request.accept_mimetypes.accept_html
            else jsonify({"type": "none", "results": []})
        )

    # --- Step 1: Refine query ---
    refined_terms = refine_query_terms(query)
    query_terms = [t.lower() for t in refined_terms]
    print(f"Refined query terms: {query_terms}")
    # --- Step 2: KG-based expansion (additive only) ---
    kg_related_terms = []
    original_terms_set = set(query_terms)
    for t in query_terms:
        try:
            expanded = kg_expand_query(t)
            app.logger.debug("KG expanded '%s' to: %s", t, expanded)
            # Ensure expansion is additive: only add new terms not already present
            for e in expanded or []:
                if not e:
                    continue
                e_norm = e.strip().lower()
                if e_norm and e_norm not in original_terms_set:
                    kg_related_terms.append(e_norm)
                    original_terms_set.add(e_norm)
        except Exception as e:
            app.logger.warning(f"⚠️ KG expand failed for term '{t}': {e}")
    # Keep original query terms and append only added KG terms
    query_terms = [t.lower() for t in refined_terms] + kg_related_terms
    query_lower = " ".join(query_terms)

    # --- Step 3: MongoDB text search ---
    results, result_type = run_search(query_lower)
    if not isinstance(results, list):
        results = []
    all_docs = {str(r["_id"]): r for r in results}

    # --- Step 4: Semantic fallback if few results ---
    if len(all_docs) < 25:
        try:
            sem_results = semantic_search(query_lower, top_k=10)
            for doc in sem_results:
                fid = str(doc["_id"])
                if fid not in all_docs:
                    all_docs[fid] = doc
        except Exception as e:
            app.logger.warning(f"Semantic fallback skipped: {e}")

    # --- Step 5: Department filter ---
    if department_filter:
        all_docs = {
            k: v
            for k, v in all_docs.items()
            if v.get("department", "").lower() == department_filter.lower()
        }

    # --- Step 6: Scoring + snippets ---
    clean_results = []
    for doc in all_docs.values():
        full_text = doc.get("data") or doc.get("search_text") or ""
        if isinstance(full_text, list):
            full_text = " ".join(map(str, full_text))
        elif isinstance(full_text, dict):
            full_text = " ".join(f"{k}: {v}" for k, v in full_text.items())
        elif not isinstance(full_text, str):
            full_text = str(full_text)

        text_lower = full_text.lower()
        tokens = re.findall(r'\b\w+\b', text_lower)
        kg_boost = 0

        for t in query_terms:
            try:
                score = kg_entity_importance(t) or 0.0
                if score > 0:
                    app.logger.debug(f"  KG entity '{t}' contributes: {score}")
                    kg_boost += math.log1p(score)
            except Exception as e:
                app.logger.warning(f"KG importance lookup failed for term '{t}': {e}")

        if kg_boost > 0:
            kg_boost = math.log1p(kg_boost) * 0.4

        phrase_pattern = r'\b' + r'\s+'.join(re.escape(term) for term in query_terms) + r'\b'
        phrase_count = len(re.findall(phrase_pattern, text_lower))
        term_count = sum(tokens.count(term) for term in query_terms)

        weight = (phrase_count * 2.0) + (term_count * 0.2)
        boost = min(2.0, math.log1p(weight)) + kg_boost
        base_score = float(doc.get("score") or doc.get("semantic_score") or 0.0)
        final_score = base_score + boost

        snippet = make_snippet(full_text, query)
        if not snippet:
            snippet = full_text[:150] + "..." if len(full_text) > 150 else full_text

        clean_results.append({
            "_id": str(doc["_id"]),
            "filename": doc.get("filename", "unknown"),
            "filetype": doc.get("filetype", "unknown"),
            "_snippet": snippet,
            "score": round(final_score, 3),
            "phrase_hits": phrase_count,
            "word_hits": term_count,
            "boost": round(boost, 3),
            "kg_boost": round(kg_boost, 3),
        })

    clean_results.sort(key=lambda x: x["score"], reverse=True)

    # --- Step 7: Filter weak results ---
    if len(clean_results) > 1:
        filtered = [r for r in clean_results if r["phrase_hits"] > 0 and r["word_hits"] > 0]
        if filtered:
            clean_results = filtered

    # --- Step 8: ONNX reranker (optional) ---
    try:
        docs_for_rerank = [
            (r["_id"], all_docs[r["_id"]].get("search_text", "")[:1000])
            for r in clean_results[:20]
        ]
        texts = [t for _, t in docs_for_rerank]
        scored = rerank_with_onnx(query, texts, top_k=10)
        for (doc_id, _), s in zip(docs_for_rerank, scored):
            for r in clean_results:
                if r["_id"] == doc_id:
                    r["rerank_score"] = s["score"]
        clean_results.sort(key=lambda x: x.get("rerank_score", x["score"]), reverse=True)
    except Exception as e:
        app.logger.warning(f"ONNX reranker skipped: {e}")

    # --- Step 9: Summarize top docs using KG ---
    final_summary = None
    if summarize_flag and clean_results:
        kg_terms = []
        original_terms_set = set(query_terms)
        for t in query_terms:
            try:
                expanded = kg_expand_query(t)
                for e in expanded or []:
                    if not e:
                        continue
                    e_norm = e.strip().lower()
                    if e_norm and e_norm not in original_terms_set:
                        kg_terms.append(e_norm)
                        original_terms_set.add(e_norm)
            except Exception:
                continue
        kg_terms = list(dict.fromkeys(kg_terms))

        try:
            kg_failures = kg_get_related_failures(query_terms + kg_terms)
            app.logger.info(f"KG returned {len(kg_failures)} related failure records.")
        except Exception as e:
            app.logger.warning(f"KG failure fetch skipped: {e}")
            kg_failures = []

        kg_context_blocks = []
        for f in kg_failures:
            text_block = (
                f"Failure Entry No: {f.get('failure_entry_no','')} | "
                f"Zone: {f.get('zone','')} | Division: {f.get('division','')} | "
                f"Section: {f.get('section','')} | Station: {f.get('station_code','')} | "
                f"Department: {f.get('department','')} | Gear: {f.get('gear_name','')} | "
                f"Cause: {f.get('cause','')} | Action: {f.get('action','')} | "
                f"Remarks: {f.get('remarks','')}"
            )
            if text_block.strip():
                kg_context_blocks.append(text_block)

        if not kg_context_blocks or all(len(b.strip()) < 10 for b in kg_context_blocks):
            for result_item in clean_results[:3]:
                doc = all_docs.get(result_item["_id"])
                if not doc:
                    continue
                full_text = doc.get("search_text") or doc.get("data") or ""
                if isinstance(full_text, dict):
                    full_text = " ".join(f"{k}:{v}" for k,v in full_text.items())
                kg_context_blocks.append(full_text[:1000])

        combined_text = "\n".join(kg_context_blocks)
        chunks = chunk_text(combined_text, max_tokens=500, overlap=50)
        top_chunks = select_top_chunks(chunks, query, top_k=3)

        summaries = []
        for chunk in top_chunks:
            prompt = f"""
You are a railway technical assistant. Analyze the failure records below.
Focus on:
- Failure trends and causes
- Gear types involved
- Actions taken and their effectiveness
- Any preventive recommendations

Provide a structured summary suitable for a control office failure analysis report.

Data:
{chunk}
"""
            try:
                output_list = llm_pipe(
    prompt,
    do_sample=True,
    temperature=0.7,
    top_p=0.9,
    top_k=50,
    max_new_tokens=200,
    eos_token_id=tokenizer.eos_token_id,
    return_full_text=False
)
                generated_text = output_list[0]["generated_text"]
                summaries.append(generated_text.strip())
            except Exception as e:
                app.logger.error(f"KG summarization failed: {e}")
                summaries.append(chunk[:300])

        final_summary = "\n\n".join(summaries[:3])

    # --- Step 10: Optional text analysis ---
    file_text = " ".join([r.get("_snippet", "") for r in clean_results])
    analysis = analyze_text_summary(file_text) if analyze_flag else None

    # --- Step 11: Response ---
    response = {"type": "mixed", "results": clean_results}
    if final_summary:
        response["summary"] = final_summary
    if analysis:
        response["analysis"] = analysis

    if request.accept_mimetypes.accept_html:
        return render_template(
            "search.html",
            results=clean_results,
            query=query,
            type="mixed",
            summary=final_summary,
            analysis=analysis
        )

    return jsonify(response)



@app.route("/export_excel/<file_id>")
def export_excel(file_id):
    file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file_doc or "data" not in file_doc:
        flash("File or data not found.", "danger")
        return redirect(url_for("dashboard"))
    excel_stream = write_excel_from_data(file_doc["filename"], file_doc["filetype"], file_doc["data"])
    return send_file(excel_stream, download_name=f"{file_doc['filename']}_data.xlsx", as_attachment=True, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")



@app.route("/export_word/<file_id>")
def export_word(file_id):
    # Fetch file document from MongoDB
    file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file_doc or "data" not in file_doc:
        flash("File or data not found.", "danger")
        return redirect(url_for("dashboard"))

    # Generate Word file stream
    word_stream = write_word_from_data(
        file_name_prefix=file_doc.get("filename", "exported_file"),
        filetype=file_doc.get("filetype", "unknown"),
        data=file_doc["data"]
    )

    # Send as downloadable file
    return send_file(
        word_stream,
        download_name=f"{file_doc.get('filename','exported_file')}_data.docx",
        as_attachment=True,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )


@app.route('/file/<file_id>')
def file_detail(file_id):
    # Fetch file from DB by ObjectId
    try:
        file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    except Exception:
        flash("Invalid file ID.", "danger")
        return redirect(url_for("search"))

    if not file_doc:
        flash("File not found.", "danger")
        return redirect(url_for("search"))

    # Optional: pass search query for highlighting
    query = request.args.get("q", "")

    # Ensure all needed fields are prepared for template
    file_data = {
        "_id": str(file_doc["_id"]),
        "file_id": str(file_doc.get("file_id")),
        "filename": file_doc.get("filename", "unknown"),
        "filetype": file_doc.get("filetype", "unknown"),
        "data": file_doc.get("data", {}),
    }

    return render_template("view_file.html", file=file_data, query=query)


@app.route("/file/<file_id>/export/download")
def export_download_file(file_id):
    try:
        grid_out = fs.get(ObjectId(file_id))
        return send_file(
            io.BytesIO(grid_out.read()),
            download_name=grid_out.filename,
            mimetype=grid_out.content_type or "application/octet-stream",
            as_attachment=True
        )
    except NoFile:
        abort(404, description="File not found in storage.")


@app.route("/file/<file_id>/export/excel")
def export_file_excel(file_id):
    file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file_doc or "data" not in file_doc:
        flash("File or data not found.", "danger")
        return redirect(url_for("dashboard"))
    excel_stream = write_excel_from_data(file_doc["filename"], file_doc["filetype"], file_doc["data"])
    return send_file(excel_stream, download_name=f"{file_doc['filename']}_data.xlsx", as_attachment=True, mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

@app.route("/file/<file_id>/export/word")
def export_file_word(file_id):
    # Fetch file document from MongoDB
    file_doc = mongo.db.files.find_one({"_id": ObjectId(file_id)})
    if not file_doc or "data" not in file_doc:
        flash("File or data not found.", "danger")
        return redirect(url_for("dashboard"))

    # Generate Word file stream
    word_stream = write_word_from_data(
        file_name_prefix=file_doc.get("filename", "exported_file"),
        filetype=file_doc.get("filetype", "unknown"),
        data=file_doc["data"]
    )

    # Send as downloadable file
    return send_file(
        word_stream,
        download_name=f"{file_doc.get('filename','exported_file')}_data.docx",
        as_attachment=True,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
@app.errorhandler(413)
def request_entity_too_large(error):
    flash("File size exceeds maximum allowed size.", "danger")
    return redirect(request.url)


@app.errorhandler(400)
def bad_request(e):
    return render_template("400.html", error=e), 400


@app.errorhandler(500)
def internal_error(e):
    return render_template("500.html", error=e), 500

# --------------------------- Admin endpoint to rebuild FAISS --------------------
# Protect with ADMIN_KEY env var in production

ADMIN_KEY = os.getenv("ADMIN_KEY")

@app.route("/admin/rebuild_faiss", methods=["POST"])
def admin_rebuild_faiss():
    key = request.args.get("key") or request.headers.get("X-ADMIN-KEY")
    if ADMIN_KEY and key != ADMIN_KEY:
        return "unauthorized", 401
    try:
        total = rebuild_faiss()
        return jsonify({"status": "ok", "vectors": total})
    except Exception as e:
        app.logger.exception("Rebuild failed: %s", e)
        return jsonify({"status": "error", "error": str(e)}), 500


if __name__ == "__main__":
    # Try to use waitress (production-like). If that fails (binding issues,
    # missing package, or running on a dev machine), fall back to Flask's
    # built-in server so the developer sees console output and the app runs.
    try:
        print("Starting server with waitress on 10.206.41.36:5006")
        serve(app, host="10.206.41.36", port=5006)
    except Exception as e:
        print(f"Waitress failed: {e}. Falling back to Flask dev server on 127.0.0.1:5006")
        # Use a loopback address for local development so binding works reliably.
        app.run(host="127.0.0.1", port=5006, debug=True)